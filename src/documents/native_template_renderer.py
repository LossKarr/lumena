"""Deterministic renderers for user-supplied Office templates."""
from __future__ import annotations

from pathlib import Path
import re
from typing import Any, Iterable

from .template_models import TemplateValidationError


_PLACEHOLDER_RE = re.compile(
    r"\{\{\s*([A-Za-z][A-Za-z0-9_. -]{0,79})\s*\}\}|"
    r"\[\[\s*([A-Za-z][A-Za-z0-9_. -]{0,79})\s*\]\]|"
    r"<<\s*([A-Za-z][A-Za-z0-9_. -]{0,79})\s*>>|"
    r"\$\{\s*([A-Za-z][A-Za-z0-9_. -]{0,79})\s*\}"
)


def _field_id(label: str) -> str:
    return re.sub(r"[^a-z0-9_]+", "_", label.lower()).strip("_") or "champ"


def _value(data: dict[str, Any], label: str) -> Any:
    normalized = _field_id(label)
    if normalized in data:
        return data[normalized]
    current: Any = data
    for part in label.replace(" ", "_").split("."):
        if not isinstance(current, dict) or part not in current:
            raise TemplateValidationError(f"missing render field: {normalized}")
        current = current[part]
    return current


def replace_placeholders(value: str, data: dict[str, Any]) -> str:
    """Replace supported placeholder syntaxes without evaluating expressions."""
    if not isinstance(value, str) or not _PLACEHOLDER_RE.search(value):
        return value

    def replacement(match: re.Match[str]) -> str:
        label = next(group for group in match.groups() if group is not None).strip()
        resolved = _value(data, label)
        return "" if resolved is None else str(resolved)

    return _PLACEHOLDER_RE.sub(replacement, value)


def _replace_runs(paragraph: Any, data: dict[str, Any]) -> None:
    runs = list(getattr(paragraph, "runs", ()))
    original = "".join(run.text for run in runs) if runs else str(getattr(paragraph, "text", ""))
    rendered = replace_placeholders(original, data)
    if rendered == original:
        return
    if runs:
        runs[0].text = rendered
        for run in runs[1:]:
            run.text = ""
    else:
        paragraph.text = rendered


def _docx_paragraphs(document: Any) -> Iterable[Any]:
    yield from document.paragraphs
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                yield from cell.paragraphs
    for section in document.sections:
        for container in (section.header, section.footer):
            yield from container.paragraphs
            for table in container.tables:
                for row in table.rows:
                    for cell in row.cells:
                        yield from cell.paragraphs


class NativeTemplateRenderer:
    """Render DOCX/XLSX/PPTX while preserving their native containers."""

    def render(self, source: Path, output: Path, data: dict[str, Any], *, renderer: str) -> Path:
        if not isinstance(data, dict):
            raise TemplateValidationError("render data must be an object")
        source_path = Path(source)
        target = Path(output)
        target.parent.mkdir(parents=True, exist_ok=True)
        if renderer == "docx-native":
            self._docx(source_path, target, data)
        elif renderer == "xlsx-native":
            self._xlsx(source_path, target, data)
        elif renderer == "pptx-native":
            self._pptx(source_path, target, data)
        else:
            raise TemplateValidationError(f"unsupported native renderer: {renderer}")
        return target

    @staticmethod
    def _docx(source: Path, output: Path, data: dict[str, Any]) -> None:
        from docx import Document

        document = Document(source)
        for paragraph in _docx_paragraphs(document):
            _replace_runs(paragraph, data)
        document.save(output)

    @staticmethod
    def _xlsx(source: Path, output: Path, data: dict[str, Any]) -> None:
        from openpyxl import load_workbook

        workbook = load_workbook(source, data_only=False)
        try:
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if isinstance(cell.value, str):
                            cell.value = replace_placeholders(cell.value, data)
            workbook.save(output)
        finally:
            workbook.close()

    @staticmethod
    def _pptx(source: Path, output: Path, data: dict[str, Any]) -> None:
        from pptx import Presentation

        presentation = Presentation(source)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for paragraph in shape.text_frame.paragraphs:
                        _replace_runs(paragraph, data)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                _replace_runs(paragraph, data)
        presentation.save(output)


__all__ = ["NativeTemplateRenderer", "replace_placeholders"]
