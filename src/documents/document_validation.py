"""Deterministic structural validation for generated document renders."""
from __future__ import annotations

from pathlib import Path
import re
from typing import Any

from PIL import Image, ImageStat


def validate_document_render(
    source: Path,
    thumbnail: Path | None,
    *,
    output_format: str,
    logo_id: str = "",
    visual_fidelity: str = "exact",
) -> dict[str, Any]:
    path = Path(source)
    fmt = str(output_format or path.suffix).lower().lstrip(".")
    proof: dict[str, Any] = {
        "status": "failed",
        "verified": False,
        "format": fmt,
        "path": str(path),
        "size": path.stat().st_size if path.is_file() else 0,
        "page_count": 0,
        "physical_page_count": 0,
        "blank_pages": [],
        "text_chars": 0,
        "thumbnail_path": str(thumbnail or ""),
        "thumbnail_width": 0,
        "thumbnail_height": 0,
        "non_blank": False,
        "logo_id": str(logo_id or ""),
        "verification_level": "visual" if visual_fidelity == "exact" else "structural",
        "errors": [],
    }
    errors: list[str] = proof["errors"]
    if not path.is_file() or proof["size"] <= 0:
        errors.append("output_missing_or_empty")
        return proof

    try:
        if fmt == "pdf":
            if not path.read_bytes()[:5].startswith(b"%PDF"):
                errors.append("invalid_pdf_header")
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            proof["physical_page_count"] = len(reader.pages)
            page_texts = [
                (page.extract_text() or "").strip()
                for page in reader.pages
            ]
            proof["blank_pages"] = [
                index
                for index, text in enumerate(page_texts, start=1)
                if len(text) < 10
            ]
            proof["page_count"] = (
                proof["physical_page_count"] - len(proof["blank_pages"])
            )
            text = "".join(page_texts)
            proof["text_chars"] = len(text.strip())
            if proof["physical_page_count"] < 1:
                errors.append("no_pages")
            if proof["blank_pages"]:
                pages = ",".join(str(index) for index in proof["blank_pages"])
                errors.append(f"blank_pdf_pages:{pages}")
            if proof["text_chars"] < 10:
                errors.append("insufficient_extractable_text")
        elif fmt == "html":
            html = path.read_text(encoding="utf-8", errors="replace")
            text = re.sub(r"<[^>]+>", " ", html)
            proof["page_count"] = 1
            proof["text_chars"] = len(re.sub(r"\s+", " ", text).strip())
            if proof["text_chars"] < 10:
                errors.append("insufficient_visible_text")
        elif fmt == "docx":
            from docx import Document

            document = Document(path)
            text = " ".join(p.text for p in document.paragraphs)
            for table in document.tables:
                text += " " + " ".join(cell.text for row in table.rows for cell in row.cells)
            proof["page_count"] = max(1, len(document.sections))
            proof["text_chars"] = len(text.strip())
            if proof["text_chars"] < 1:
                errors.append("insufficient_extractable_text")
        elif fmt == "xlsx":
            from openpyxl import load_workbook

            workbook = load_workbook(path, data_only=False, read_only=True)
            try:
                values = [
                    str(cell.value)
                    for sheet in workbook.worksheets
                    for row in sheet.iter_rows()
                    for cell in row
                    if cell.value not in (None, "")
                ]
                proof["page_count"] = len(workbook.worksheets)
                proof["text_chars"] = len(" ".join(values).strip())
                if proof["page_count"] < 1:
                    errors.append("no_worksheets")
                if proof["text_chars"] < 1:
                    errors.append("insufficient_extractable_text")
            finally:
                workbook.close()
        elif fmt == "pptx":
            from pptx import Presentation

            presentation = Presentation(path)
            values = [
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            proof["page_count"] = len(presentation.slides)
            proof["text_chars"] = len(" ".join(values).strip())
            if proof["page_count"] < 1:
                errors.append("no_slides")
            if proof["text_chars"] < 1:
                errors.append("insufficient_extractable_text")
        else:
            errors.append("unsupported_format")
    except Exception as exc:
        errors.append(f"content_validation_failed:{type(exc).__name__}")

    thumb = Path(thumbnail) if thumbnail else None
    thumbnail_required = visual_fidelity == "exact"
    if thumb is None or not thumb.is_file():
        if thumbnail_required:
            errors.append("thumbnail_missing")
    else:
        try:
            with Image.open(thumb) as image:
                rgb = image.convert("RGB")
                proof["thumbnail_width"], proof["thumbnail_height"] = rgb.size
                gray = rgb.convert("L")
                extrema = gray.getextrema()
                variance = float(ImageStat.Stat(gray).var[0])
                proof["non_blank"] = bool(extrema and extrema[1] - extrema[0] >= 8 and variance >= 1.0)
                if not proof["non_blank"]:
                    errors.append("thumbnail_blank")
        except Exception as exc:
            errors.append(f"thumbnail_validation_failed:{type(exc).__name__}")

    proof["verified"] = not errors
    if proof["verified"]:
        proof["status"] = "render_verified" if visual_fidelity == "exact" else "structure_verified"
    else:
        proof["status"] = "render_failed"
    return proof


__all__ = ["validate_document_render"]
