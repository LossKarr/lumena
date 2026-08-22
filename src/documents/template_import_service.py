"""Safe conversion of user documents into reviewable HTML/Jinja template drafts."""
from __future__ import annotations

from dataclasses import replace
from html import escape
import json
from pathlib import Path
import re
import shutil
from typing import Any
from uuid import uuid4

from .document_security import inspect_document, sanitize_document_filename
from .import_service import DocumentImportService
from .template_import_models import TemplateImportDraft
from .template_security import atomic_write_json, normalize_template_id, resolve_within, validate_template_source


_PLACEHOLDER = re.compile(
    r"\{\{\s*([A-Za-z][A-Za-z0-9_. -]{0,79})\s*\}\}|"
    r"\[\[\s*([A-Za-z][A-Za-z0-9_. -]{0,79})\s*\]\]|"
    r"<<\s*([A-Za-z][A-Za-z0-9_. -]{0,79})\s*>>|"
    r"\$\{\s*([A-Za-z][A-Za-z0-9_. -]{0,79})\s*\}"
)
_FIELD_ID = re.compile(r"[^a-z0-9_]+")


class TemplateImportService:
    """Store immutable sources and mutable drafts; publication is always explicit."""

    def __init__(self, root: Path, catalog):
        self.root = Path(root)
        self.catalog = catalog
        self.root.mkdir(parents=True, exist_ok=True)

    def list_drafts(self) -> list[dict[str, Any]]:
        drafts = []
        for directory in sorted(self.root.glob("draft_*")):
            try:
                drafts.append(self._load(directory.name).to_dict(include_source=False))
            except Exception:
                continue
        return drafts

    def get(self, draft_id: str) -> TemplateImportDraft:
        return self._load(draft_id)

    def create(self, source: Path, *, filename: str = "", name: str = "",
               kind: str = "", category: str = "custom") -> TemplateImportDraft:
        report = inspect_document(Path(source))
        if "external_relationships" in report.warnings:
            raise ValueError("Les documents avec relations externes ne peuvent pas devenir un modèle")
        draft_id = f"draft_{uuid4().hex[:16]}"
        directory = resolve_within(self.root, draft_id)
        directory.mkdir(parents=True, exist_ok=False)
        safe_name = sanitize_document_filename(filename or Path(source).name)
        source_name = f"source{Path(safe_name).suffix.lower()}"
        shutil.copyfile(source, directory / source_name)
        try:
            extracted = self._extract(directory / source_name, report)
            display_name = str(name or Path(safe_name).stem.replace("_", " ")).strip()
            template_kind = normalize_template_id(kind or Path(safe_name).stem)
            draft = TemplateImportDraft.create(
                draft_id=draft_id, source_filename=safe_name,
                source_format=report.format, source_sha256=report.sha256,
                source_file=source_name, name=display_name, kind=template_kind,
                template_source=extracted["source"], sample_data=extracted["sample_data"],
                detected_fields=extracted["fields"], fidelity=extracted["fidelity"],
                warnings=[*report.warnings, *extracted["warnings"]],
            )
            native_renderer = {
                "docx": "docx-native",
                "xlsx": "xlsx-native",
                "pptx": "pptx-native",
            }.get(report.format)
            draft = replace(
                draft,
                category=str(category or "custom").strip().lower(),
                renderer=native_renderer or "html-jinja",
                output_format=report.format if native_renderer else "pdf",
            )
            self._write(draft)
            return draft
        except Exception:
            shutil.rmtree(directory, ignore_errors=True)
            raise

    def update(self, draft_id: str, payload: dict[str, Any]) -> TemplateImportDraft:
        current = self._load(draft_id)
        if current.status != "draft":
            raise ValueError("Un modèle publié ne peut plus modifier son brouillon source")
        aliases = payload.get("aliases", current.aliases)
        if not isinstance(aliases, (list, tuple)):
            raise ValueError("aliases doit être une liste")
        sample = payload.get("sample_data", current.sample_data)
        if not isinstance(sample, dict):
            raise ValueError("sample_data doit être un objet")
        source = str(payload.get("template_source", current.template_source))
        validate_template_source(source)
        from datetime import datetime, timezone
        updated = replace(
            current,
            updated_at=datetime.now(timezone.utc).isoformat(),
            name=str(payload.get("name", current.name)).strip(),
            kind=normalize_template_id(str(payload.get("kind", current.kind))),
            category=normalize_template_id(str(payload.get("category", current.category))),
            aliases=tuple(str(value).strip() for value in aliases if str(value).strip()),
            template_source=source,
            sample_data=dict(sample),
            detected_fields=tuple(dict(value) for value in payload.get("detected_fields", current.detected_fields)),
        )
        updated.validate()
        self._write(updated)
        return updated

    def publish(self, draft_id: str, *, template_id: str = ""):
        draft = self._load(draft_id)
        if draft.status != "draft":
            raise ValueError("Ce brouillon est déjà publié")
        target_id = normalize_template_id(template_id or draft.kind)
        native = draft.renderer != "html-jinja"
        template_filename = f"template.{draft.output_format}" if native else "template.html.j2"
        manifest = {
            "schema_version": 1, "id": target_id, "name": draft.name,
            "kind": draft.kind, "format": draft.output_format,
            "renderer": draft.renderer, "origin": "custom", "version": 1,
            "editable_fields": list(draft.detected_fields), "design": {},
            "sample_data": "sample-data.json", "template_file": template_filename,
            "thumbnail": "thumbnail.webp", "description": f"Modèle importé depuis {draft.source_filename}",
            "category": draft.category, "aliases": list(draft.aliases),
            "locale": "fr-FR", "scope": "universal", "compliance_level": "structure",
            "jurisdictions": [], "legal_notice": "",
        }
        if native:
            record = self.catalog.save_custom_native(
                target_id,
                manifest_data=manifest,
                source_file=resolve_within(self.root, draft.id) / draft.source_file,
                sample_data=draft.sample_data,
            )
        else:
            record = self.catalog.save_custom(
                target_id, manifest_data=manifest,
                template_source=draft.template_source, sample_data=draft.sample_data,
            )
        published = replace(draft, status="published", published_template_id=record.manifest.id)
        self._write(published)
        return record

    def preview_html(self, draft_id: str) -> str:
        from .template_models import TemplateManifest, TemplateRecord
        from .template_renderer import TemplateRenderer

        draft = self._load(draft_id)
        manifest = TemplateManifest(
            schema_version=1, id="draft-preview", name=draft.name, kind="draft_preview",
            format="pdf", renderer="html-jinja", origin="custom",
        )
        record = TemplateRecord(manifest, resolve_within(self.root, draft.id), False)
        return TemplateRenderer().render_html(
            record, draft.template_source, draft.sample_data,
        )

    def delete(self, draft_id: str) -> None:
        draft = self._load(draft_id)
        if draft.status == "published":
            raise ValueError("Le brouillon publié est conservé comme preuve de provenance")
        shutil.rmtree(resolve_within(self.root, draft_id))

    def _load(self, draft_id: str) -> TemplateImportDraft:
        path = resolve_within(self.root, draft_id) / "draft.json"
        if not path.is_file():
            raise KeyError(draft_id)
        return TemplateImportDraft.from_dict(json.loads(path.read_text(encoding="utf-8")))

    def _write(self, draft: TemplateImportDraft) -> None:
        draft.validate()
        atomic_write_json(resolve_within(self.root, draft.id) / "draft.json", draft.to_dict())

    def _extract(self, path: Path, report) -> dict[str, Any]:
        if report.format == "html":
            body = self._safe_html(path.read_text(encoding="utf-8", errors="replace"))
            fidelity = "high"
        elif report.format == "docx":
            body = self._docx_html(path)
            fidelity = "structural"
        elif report.format == "xlsx":
            body = self._xlsx_html(path)
            fidelity = "data-and-style-approximation"
        elif report.format == "pptx":
            body = self._pptx_html(path)
            fidelity = "structural"
        else:
            text = DocumentImportService._extract_text(path, report)
            body = "<article><h1>" + escape(path.stem.replace("_", " ").title()) + "</h1>" + "".join(
                f"<p>{escape(line)}</p>" for line in text.splitlines() if line.strip()
            ) + "</article>"
            fidelity = "semantic" if text.strip() else "reference-only"
        source, fields, sample = self._normalize_placeholders(body)
        wrapped = self._wrap(source)
        validate_template_source(wrapped)
        warnings = [] if fields else ["no_placeholders_detected"]
        return {"source": wrapped, "fields": fields, "sample_data": sample,
                "fidelity": fidelity, "warnings": warnings}

    @classmethod
    def structural_preview_html(cls, path: Path, source_format: str) -> str:
        """Create an explicitly structural, non-fidelity preview for Office files."""
        fmt = str(source_format).lower().lstrip(".")
        if fmt == "docx":
            body = cls._docx_html(Path(path))
        elif fmt == "xlsx":
            body = cls._xlsx_html(Path(path))
        elif fmt == "pptx":
            body = cls._pptx_html(Path(path))
        else:
            raise ValueError(f"Aperçu structurel non pris en charge pour {fmt}")
        return cls._wrap(body)

    @staticmethod
    def _safe_html(raw: str) -> str:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup.find_all(["script", "iframe", "object", "embed", "form", "input", "button", "base"]):
            tag.decompose()
        for tag in soup.find_all(True):
            for attr in list(tag.attrs):
                value = str(tag.attrs.get(attr, ""))
                if attr.lower().startswith("on") or "javascript:" in value.lower():
                    del tag.attrs[attr]
                elif attr.lower() in {"src", "href"} and value.lower().startswith(("http://", "https://", "//")):
                    del tag.attrs[attr]
            if tag.name == "style" and tag.string:
                tag.string.replace_with(re.sub(r"@import[^;]+;|url\([^)]*\)", "", tag.string, flags=re.I))
        return str(soup.body or soup)

    @staticmethod
    def _docx_html(path: Path) -> str:
        from docx import Document
        document = Document(path)
        parts = []
        for paragraph in document.paragraphs:
            text = escape(paragraph.text)
            if not text:
                continue
            style = str(paragraph.style.name or "").lower()
            tag = "h1" if "title" in style else "h2" if "heading" in style else "p"
            parts.append(f"<{tag}>{text}</{tag}>")
        for table in document.tables:
            parts.append("<table>")
            for row in table.rows:
                parts.append("<tr>" + "".join(f"<td>{escape(cell.text)}</td>" for cell in row.cells) + "</tr>")
            parts.append("</table>")
        return "<article>" + "".join(parts) + "</article>"

    @staticmethod
    def _xlsx_html(path: Path) -> str:
        from openpyxl import load_workbook
        workbook = load_workbook(path, data_only=False, read_only=True)
        sheet = workbook.active
        rows = []
        for row in sheet.iter_rows(max_row=min(sheet.max_row, 500), max_col=min(sheet.max_column, 50)):
            rows.append("<tr>" + "".join(f"<td>{escape(str(cell.value or ''))}</td>" for cell in row) + "</tr>")
        workbook.close()
        return f"<article><h1>{escape(sheet.title)}</h1><table>{''.join(rows)}</table></article>"

    @staticmethod
    def _pptx_html(path: Path) -> str:
        from pptx import Presentation
        presentation = Presentation(path)
        slides = []
        for index, slide in enumerate(presentation.slides, 1):
            texts = [escape(shape.text) for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            slides.append(f"<section class='imported-slide'><h2>Page {index}</h2>" + "".join(f"<p>{text}</p>" for text in texts) + "</section>")
        return "<article>" + "".join(slides) + "</article>"

    @classmethod
    def _normalize_placeholders(cls, html: str):
        fields: list[dict[str, Any]] = []
        sample: dict[str, Any] = {}
        seen: set[str] = set()

        def replace_match(match):
            label = next(value for value in match.groups() if value is not None).strip()
            field_id = _FIELD_ID.sub("_", label.lower()).strip("_") or "champ"
            if field_id not in seen:
                seen.add(field_id)
                fields.append({"id": field_id, "label": label, "type": "text", "required": False})
                sample[field_id] = label
            return "{{ " + field_id + " }}"

        return _PLACEHOLDER.sub(replace_match, html), fields, sample

    @staticmethod
    def _wrap(body: str) -> str:
        if "<!doctype" in body.lower():
            return body
        return """<!DOCTYPE html><html lang=\"fr\"><head><meta charset=\"UTF-8\"><style>
@page{size:A4;margin:18mm}body{font-family:Arial,sans-serif;color:#1c2430;font-size:11px;line-height:1.5}
article{max-width:100%}h1,h2{color:#b45309}table{width:100%;border-collapse:collapse;margin:14px 0}
td,th{border:1px solid #d0d5dd;padding:8px;vertical-align:top}.imported-slide{break-after:page}
</style></head><body>""" + body + "</body></html>"
