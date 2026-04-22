"""
Document Hub — Génération PDF, DOCX, XLSX, PPTX pour Lumena.

Supporte le markdown inline (gras, italique, titres, listes, hr, tables).
Tous les fichiers sont sauvegardés dans workspace/YYYY-MM-DD/.
"""
from __future__ import annotations

import csv
import io
import json
import os
import re
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

from loguru import logger


# ──────────────────────────────────────────────────────────
# Page-size / orientation helpers
# ──────────────────────────────────────────────────────────

_PAGE_SIZES = {}  # populated lazily

def _get_page_size(name: str = "A4", orientation: str = "portrait"):
    """Return reportlab page-size tuple.  Lazy-import."""
    global _PAGE_SIZES
    if not _PAGE_SIZES:
        from reportlab.lib.pagesizes import A3, A4, A5, legal, letter
        _PAGE_SIZES = {"A4": A4, "A3": A3, "A5": A5, "LETTER": letter, "LEGAL": legal}
    ps = _PAGE_SIZES.get(name.upper(), _PAGE_SIZES["A4"])
    if orientation.lower() == "landscape":
        from reportlab.lib.pagesizes import landscape
        ps = landscape(ps)
    return ps


# ──────────────────────────────────────────────────────────
# Markdown table parser (shared by PDF + DOCX)
# ──────────────────────────────────────────────────────────

def _parse_md_tables(content: str):
    """Parse markdown tables in content.

    Returns a list of (start_line_idx, end_line_idx, headers, rows) and
    the lines list for reference.
    """
    lines = content.split("\n")
    tables: list = []
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        if line.startswith("|") and line.endswith("|"):
            # Potential table start
            header_cells = [c.strip() for c in line.strip("|").split("|")]
            start = i
            i += 1
            # skip separator |---|---|
            if i < len(lines) and re.match(r"^\|[\s\-:|]+\|$", lines[i].strip()):
                i += 1
            rows = []
            while i < len(lines) and lines[i].strip().startswith("|") and lines[i].strip().endswith("|"):
                row_cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                rows.append(row_cells)
                i += 1
            tables.append((start, i, header_cells, rows))
        else:
            i += 1
    return lines, tables


class DocumentHub:
    """Génère des documents Office et PDF depuis du contenu structuré."""

    # ── P4.1 — Thèmes prédéfinis ──
    _DOCUMENT_THEMES: Dict[str, Dict[str, Any]] = {
        "corporate": {
            "primary": "#000000", "secondary": "#000000", "text": "#000000",
            "bg": "#ffffff", "accent": "#000000",
            "font_heading": "Helvetica-Bold", "font_body": "Helvetica",
            "font_size_body": 11, "font_size_h1": 22, "font_size_h2": 16, "font_size_h3": 13,
        },
        "minimal": {
            "primary": "#000000", "secondary": "#000000", "text": "#000000",
            "bg": "#ffffff", "accent": "#000000",
            "font_heading": "Helvetica-Bold", "font_body": "Helvetica",
            "font_size_body": 11, "font_size_h1": 20, "font_size_h2": 15, "font_size_h3": 12,
        },
        "modern": {
            "primary": "#6366f1", "secondary": "#22c55e", "text": "#1e293b",
            "bg": "#f8fafc", "accent": "#6366f1",
            "font_heading": "Helvetica-Bold", "font_body": "Helvetica",
            "font_size_body": 11, "font_size_h1": 24, "font_size_h2": 18, "font_size_h3": 14,
        },
        "legal": {
            "primary": "#000000", "secondary": "#000000", "text": "#000000",
            "bg": "#ffffff", "accent": "#000000",
            "font_heading": "Times-Bold", "font_body": "Times-Roman",
            "font_size_body": 12, "font_size_h1": 18, "font_size_h2": 15, "font_size_h3": 13,
        },
        "creative": {
            "primary": "#e11d48", "secondary": "#f59e0b", "text": "#1e293b",
            "bg": "#fffbeb", "accent": "#e11d48",
            "font_heading": "Helvetica-Bold", "font_body": "Helvetica",
            "font_size_body": 11, "font_size_h1": 26, "font_size_h2": 18, "font_size_h3": 14,
        },
    }

    def __init__(self, workspace_root: Path):
        self.workspace_root = Path(workspace_root)
        self.workspace_root.mkdir(parents=True, exist_ok=True)

    # ──────────────────────────────────────────────────────
    # Utilitaires internes
    # ──────────────────────────────────────────────────────

    def _output_path(self, filename: str) -> Path:
        """
        Chemin de sortie.
        - Si filename est un chemin absolu → utilisé directement (overwrite possible).
        - Sinon → workspace/YYYY-MM-DD/filename.
        """
        p = Path(filename)
        if p.is_absolute():
            p.parent.mkdir(parents=True, exist_ok=True)
            return p
        today = datetime.now().strftime("%Y-%m-%d")
        out_dir = self.workspace_root / today
        out_dir.mkdir(parents=True, exist_ok=True)
        return out_dir / filename

    @staticmethod
    def _ensure_ext(filename: str, ext: str) -> str:
        if not filename.lower().endswith(ext):
            filename += ext
        return filename

    @staticmethod
    def _parse_json_arg(arg: Any, fallback: Any = None) -> Any:
        """Parse JSON si c'est une string, sinon retourne tel quel."""
        if isinstance(arg, str):
            try:
                return json.loads(arg)
            except Exception:
                return fallback
        return arg if arg is not None else fallback

    def _get_theme(self, theme: Optional[str] = None) -> Dict[str, Any]:
        """Retourne le thème demandé, l'env var, ou le thème corporate par défaut."""
        name = theme or os.getenv("LUMENA_DOCUMENT_THEME", "")
        return self._DOCUMENT_THEMES.get(name, self._DOCUMENT_THEMES["corporate"])

    @staticmethod
    def _format_page_number(page_num: int, numbering: str = "arabic") -> str:
        """Formate un numéro de page selon le style demandé."""
        if numbering == "roman":
            vals = [(1000, "m"), (900, "cm"), (500, "d"), (400, "cd"),
                    (100, "c"), (90, "xc"), (50, "l"), (40, "xl"),
                    (10, "x"), (9, "ix"), (5, "v"), (4, "iv"), (1, "i")]
            result = ""
            for v, r in vals:
                while page_num >= v:
                    result += r
                    page_num -= v
            return result
        elif numbering == "Roman":
            return DocumentHub._format_page_number(page_num, "roman").upper()
        elif numbering == "alpha":
            if 1 <= page_num <= 26:
                return chr(96 + page_num)
            return str(page_num)
        elif numbering == "Alpha":
            if 1 <= page_num <= 26:
                return chr(64 + page_num)
            return str(page_num)
        return str(page_num)

    @staticmethod
    def _make_inline_chart(chart_def: Dict) -> Optional[str]:
        """Génère un chart matplotlib en PNG temp, retourne le chemin ou None."""
        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt

            chart_type = chart_def.get("chart_type", "bar")
            data = chart_def.get("data", {})
            title = chart_def.get("title", "")
            labels = data.get("labels", [])
            values = data.get("values", [])

            if not labels or not values:
                return None

            fig, ax = plt.subplots(figsize=(8, 4.5))
            if chart_type == "bar":
                ax.bar(labels, values, color="#6366f1")
            elif chart_type == "line":
                ax.plot(labels, values, marker="o", color="#6366f1")
            elif chart_type == "pie":
                ax.pie(values, labels=labels, autopct="%1.1f%%")
            elif chart_type == "scatter":
                ax.scatter(labels, values, color="#6366f1")
            else:
                ax.bar(labels, values, color="#6366f1")

            if title:
                ax.set_title(title)
            fig.tight_layout()
            import tempfile
            fd, tmp_path = tempfile.mkstemp(suffix=".png")
            os.close(fd)
            fig.savefig(tmp_path, dpi=150, bbox_inches="tight")
            plt.close(fig)
            return tmp_path
        except Exception as e:
            logger.warning(f"Impossible de générer le chart inline: {e}")
            return None

    @staticmethod
    def _parse_footnotes(content: str) -> tuple:
        """Parse les footnotes [^N]: ... du contenu. Retourne (contenu_nettoyé, dict_footnotes)."""
        footnotes: Dict[str, str] = {}
        lines = content.split("\n")
        clean_lines = []
        for line in lines:
            m = re.match(r"^\[\^(\w+)\]:\s*(.*)", line.strip())
            if m:
                footnotes[m.group(1)] = m.group(2)
            else:
                clean_lines.append(line)
        return "\n".join(clean_lines), footnotes

    # ──────────────────────────────────────────────────────
    # PDF (via reportlab — déjà installé)
    # ──────────────────────────────────────────────────────

    @staticmethod
    def _md_inline_rl(text: str) -> str:
        """Convertit markdown inline → balises ReportLab XML."""
        # Extraire les zones markdown AVANT d'échapper, pour préserver les marqueurs
        # Étape 1 : échapper les caractères XML spéciaux bruts (hors marqueurs markdown)
        # On protège d'abord les séquences markdown, on échappe, puis on restaure
        import html as _html
        # Reconstruire proprement : échapper & < > sauf dans les balises qu'on va créer
        # Approche simple : échapper tout, puis appliquer les substitutions markdown
        text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        # Bold italic
        text = re.sub(r"\*\*\*(.*?)\*\*\*", r"<b><i>\1</i></b>", text)
        # Bold
        text = re.sub(r"\*\*(.*?)\*\*", r"<b>\1</b>", text)
        # Italic
        text = re.sub(r"\*(.*?)\*", r"<i>\1</i>", text)
        # Inline code
        text = re.sub(r"`(.*?)`", r"<font name='Courier'>\1</font>", text)
        return text

    def create_pdf(
        self,
        filename: str,
        title: str,
        content: str,
        font_size: int = 11,
        author: str = "Lumena",
        images: Optional[List[Dict]] = None,
        header_footer: Optional[Dict] = None,
        page_size: str = "A4",
        orientation: str = "portrait",
        columns: int = 1,
        numbering: str = "arabic",
        toc: bool = False,
        charts: Optional[List[Dict]] = None,
        theme: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Crée un PDF depuis du contenu texte/markdown.
        Supporte : # ## ### titres, **gras**, *italique*, `code`, listes -, 1.,
        séparateurs ---, tables markdown |...|, images, headers/footers,
        thèmes, TOC, checkboxes, footnotes, multi-colonnes, numérotation.
        """
        try:
            from reportlab.lib import colors
            from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT, TA_RIGHT
            from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
            from reportlab.lib.units import cm
            from reportlab.platypus import (
                HRFlowable,
                Image as RLImage,
                Paragraph,
                SimpleDocTemplate,
                Spacer,
                Table,
                TableStyle,
            )

            pagesize = _get_page_size(page_size, orientation)
            images = self._parse_json_arg(images, [])
            header_footer = self._parse_json_arg(header_footer, None)
            charts = self._parse_json_arg(charts, [])

            filename = self._ensure_ext(filename, ".pdf")
            out_path = self._output_path(filename)

            # ── P4.1 Theme ──
            th = self._get_theme(theme)

            # ── P4.5 Footnotes ──
            content, footnotes_dict = self._parse_footnotes(content)
            _footnotes_collected: List[str] = []  # (label, text) collected during parse
            _numbering = numbering

            # Header/footer callback with numbering format
            _hf = header_footer

            def _on_page(canvas, doc_obj):
                canvas.saveState()
                page_w, page_h = pagesize
                if _hf:
                    hl = _hf.get("header_left", "")
                    hr_ = _hf.get("header_right", "")
                    if hl:
                        canvas.setFont(th.get("font_body", "Helvetica"), 8)
                        canvas.drawString(2 * cm, page_h - 1.2 * cm, hl)
                    if hr_:
                        canvas.setFont(th.get("font_body", "Helvetica"), 8)
                        canvas.drawRightString(page_w - 2 * cm, page_h - 1.2 * cm, hr_)
                    fc = _hf.get("footer_center", "")
                    if fc:
                        page_num = canvas.getPageNumber()
                        formatted = DocumentHub._format_page_number(page_num, _numbering)
                        text = fc.replace("{page}", formatted).replace("{total}", "?")
                        canvas.setFont(th.get("font_body", "Helvetica"), 8)
                        canvas.drawCentredString(page_w / 2, 1 * cm, text)
                # ── P4.5 Footnotes at bottom ──
                if _footnotes_collected:
                    y = 2 * cm
                    canvas.setFont(th.get("font_body", "Helvetica"), 7)
                    canvas.setFillColor(colors.HexColor("#666666"))
                    canvas.line(2 * cm, y + 4, page_w / 2, y + 4)
                    for i, fn_text in enumerate(_footnotes_collected):
                        canvas.drawString(2 * cm, y - (i * 10), fn_text)
                canvas.restoreState()

            # ── P4.6 Multi-column layout ──
            columns = max(1, min(columns, 3))
            if columns > 1:
                from reportlab.platypus import BaseDocTemplate, Frame, PageTemplate
                gutter = 1 * cm
                page_w, page_h = pagesize
                margin_lr = 2 * cm
                margin_top = 2.5 * cm if not _hf else 3 * cm
                margin_bot = 2 * cm if not _hf else 2.5 * cm
                usable_w = page_w - 2 * margin_lr
                col_w = (usable_w - (columns - 1) * gutter) / columns
                frames = []
                for ci in range(columns):
                    x = margin_lr + ci * (col_w + gutter)
                    frames.append(Frame(x, margin_bot, col_w, page_h - margin_top - margin_bot))
                pt = PageTemplate(id="multi", frames=frames, onPage=_on_page)
                doc = BaseDocTemplate(
                    str(out_path), pagesize=pagesize, title=title, author=author,
                )
                doc.addPageTemplates([pt])
            else:
                doc = SimpleDocTemplate(
                    str(out_path),
                    pagesize=pagesize,
                    rightMargin=2 * cm,
                    leftMargin=2 * cm,
                    topMargin=2.5 * cm if not _hf else 3 * cm,
                    bottomMargin=2 * cm if not _hf else 2.5 * cm,
                    title=title,
                    author=author,
                )

            styles = getSampleStyleSheet()

            style_title = ParagraphStyle(
                "LTitle",
                parent=styles["Title"],
                fontSize=th["font_size_h1"] + 2,
                fontName=th["font_heading"],
                textColor=colors.HexColor(th["primary"]),
                spaceAfter=16,
                spaceBefore=0,
                alignment=TA_CENTER,
            )
            style_h1 = ParagraphStyle(
                "LH1",
                parent=styles["Heading1"],
                fontSize=th["font_size_h1"],
                fontName=th["font_heading"],
                textColor=colors.HexColor(th["primary"]),
                spaceBefore=14,
                spaceAfter=6,
            )
            style_h2 = ParagraphStyle(
                "LH2",
                parent=styles["Heading2"],
                fontSize=th["font_size_h2"],
                fontName=th["font_heading"],
                textColor=colors.HexColor(th["secondary"]),
                spaceBefore=10,
                spaceAfter=4,
            )
            style_h3 = ParagraphStyle(
                "LH3",
                parent=styles["Heading3"],
                fontSize=th["font_size_h3"],
                fontName=th["font_heading"],
                textColor=colors.HexColor(th.get("text", "#4a5568")),
                spaceBefore=8,
                spaceAfter=3,
            )
            style_body = ParagraphStyle(
                "LBody",
                parent=styles["Normal"],
                fontSize=font_size,
                fontName=th["font_body"],
                leading=font_size * 1.5,
                spaceAfter=5,
                alignment=TA_JUSTIFY,
            )
            style_bullet = ParagraphStyle(
                "LBullet",
                parent=styles["Normal"],
                fontSize=font_size,
                fontName=th["font_body"],
                leading=font_size * 1.4,
                leftIndent=18,
                firstLineIndent=0,
                spaceAfter=2,
            )
            style_footnote = ParagraphStyle(
                "LFootnote",
                parent=styles["Normal"],
                fontSize=7,
                fontName=th["font_body"],
                textColor=colors.HexColor("#666666"),
                leading=9,
            )

            COLOR_HR = colors.HexColor("#aaaaaa")
            COLOR_HR_TITLE = colors.HexColor(th.get("accent", "#2d3561"))

            story = []
            _toc_entries: List[tuple] = []  # (level, text) for bookmarks

            story.append(Paragraph(title, style_title))
            story.append(
                HRFlowable(
                    width="100%",
                    thickness=1.5,
                    color=COLOR_HR_TITLE,
                    spaceAfter=14,
                )
            )

            def _apply_footnote_refs(text: str) -> str:
                """Remplace [^N] par exposant dans le texte inline."""
                def _fn_repl(m):
                    key = m.group(1)
                    if key in footnotes_dict:
                        fn_num = list(footnotes_dict.keys()).index(key) + 1
                        note_text = footnotes_dict[key]
                        _footnotes_collected.append(f"{fn_num}. {note_text}")
                        return f"<super><font size='7'>{fn_num}</font></super>"
                    return m.group(0)
                return re.sub(r"\[\^(\w+)\]", _fn_repl, text)

            # Parse tables from content
            all_lines, md_tables = _parse_md_tables(content)
            table_line_ranges = set()
            for (ts, te, _, _) in md_tables:
                table_line_ranges.update(range(ts, te))

            for line_idx, line in enumerate(all_lines):
                if line_idx in table_line_ranges:
                    for (ts, te, t_headers, t_rows) in md_tables:
                        if line_idx == ts:
                            tbl_data = [[Paragraph(f"<b>{self._md_inline_rl(h)}</b>", style_body) for h in t_headers]]
                            for row in t_rows:
                                tbl_data.append([Paragraph(self._md_inline_rl(c), style_body) for c in row])
                            n_cols = len(t_headers) if t_headers else 1
                            page_w_tbl = pagesize[0] - 4 * cm
                            col_w_tbl = page_w_tbl / max(n_cols, 1)
                            tbl = Table(tbl_data, colWidths=[col_w_tbl] * n_cols)
                            tbl.setStyle(TableStyle([
                                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(th["primary"])),
                                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                                ("FONTNAME", (0, 0), (-1, 0), th["font_heading"]),
                                ("FONTSIZE", (0, 0), (-1, -1), font_size - 1),
                                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f0f4ff")]),
                                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cccccc")),
                                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                                ("TOPPADDING", (0, 0), (-1, -1), 4),
                                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                            ]))
                            story.append(Spacer(1, 6))
                            story.append(tbl)
                            story.append(Spacer(1, 6))
                    continue

                s = line.strip()
                if not s:
                    story.append(Spacer(1, 4))
                    continue

                # Inline image
                img_match = re.match(r"^!\[([^\]]*)\]\(([^)]+)\)\s*$", s)
                if img_match:
                    caption, img_path = img_match.group(1), img_match.group(2)
                    if Path(img_path).exists():
                        story.append(Spacer(1, 4))
                        story.append(RLImage(img_path, width=12 * cm, height=8 * cm, kind="proportional"))
                        if caption:
                            story.append(Paragraph(f"<i>{self._md_inline_rl(caption)}</i>", style_body))
                        story.append(Spacer(1, 4))
                    continue

                # ── P4.4 Checkboxes ──
                cb_match = re.match(r"^- \[([ xX])\]\s+(.*)", s)
                if cb_match:
                    checked = cb_match.group(1).lower() == "x"
                    cb_char = "☑" if checked else "☐"
                    cb_text = _apply_footnote_refs(self._md_inline_rl(cb_match.group(2)))
                    story.append(Paragraph(f"{cb_char} {cb_text}", style_bullet))
                    continue

                # Apply footnote refs to inline text
                s_fn = _apply_footnote_refs(s)

                if s.startswith("### "):
                    text = self._md_inline_rl(s[4:])
                    story.append(Paragraph(_apply_footnote_refs(text), style_h3))
                    if toc:
                        _toc_entries.append((2, s[4:]))
                elif s.startswith("## "):
                    text = self._md_inline_rl(s[3:])
                    story.append(Paragraph(_apply_footnote_refs(text), style_h2))
                    if toc:
                        _toc_entries.append((1, s[3:]))
                elif s.startswith("# "):
                    text = self._md_inline_rl(s[2:])
                    story.append(Paragraph(_apply_footnote_refs(text), style_h1))
                    story.append(HRFlowable(width="100%", thickness=0.5, color=COLOR_HR, spaceAfter=4))
                    if toc:
                        _toc_entries.append((0, s[2:]))
                elif s.startswith(("- ", "* ", "• ")):
                    story.append(
                        Paragraph(f"• {_apply_footnote_refs(self._md_inline_rl(s[2:]))}", style_bullet)
                    )
                elif re.match(r"^\d+\.\s", s):
                    m = re.match(r"^(\d+)\.\s+(.*)", s)
                    if m:
                        idx = int(m.group(1))
                        formatted_idx = self._format_page_number(idx, numbering)
                        story.append(
                            Paragraph(
                                f"{formatted_idx}.  {_apply_footnote_refs(self._md_inline_rl(m.group(2)))}",
                                style_bullet,
                            )
                        )
                elif re.match(r"^-{3,}$|^={3,}$", s):
                    story.append(
                        HRFlowable(width="100%", thickness=0.5, color=COLOR_HR, spaceBefore=6, spaceAfter=6)
                    )
                else:
                    story.append(Paragraph(_apply_footnote_refs(self._md_inline_rl(s)), style_body))

            # Append explicit images
            for img_def in (images or []):
                img_path = img_def.get("path", "")
                if img_path and Path(img_path).exists():
                    w = img_def.get("width_cm", 12) * cm
                    story.append(Spacer(1, 6))
                    story.append(RLImage(img_path, width=w, height=w * 0.67, kind="proportional"))
                    cap = img_def.get("caption", "")
                    if cap:
                        story.append(Paragraph(f"<i>{self._md_inline_rl(cap)}</i>", style_body))

            # ── P4.3 Inline charts (generate from data) ──
            _chart_temps = []
            for chart_def in (charts or []):
                chart_path = chart_def.get("path", "")
                if not chart_path or not Path(chart_path).exists():
                    # Try to generate inline
                    gen_path = self._make_inline_chart(chart_def)
                    if gen_path:
                        chart_path = gen_path
                        _chart_temps.append(gen_path)
                if chart_path and Path(chart_path).exists():
                    w = chart_def.get("width_cm", 14) * cm
                    story.append(Spacer(1, 8))
                    story.append(RLImage(chart_path, width=w, height=w * 0.56, kind="proportional"))
                    cap = chart_def.get("caption", chart_def.get("title", ""))
                    if cap:
                        story.append(Paragraph(f"<i>{self._md_inline_rl(cap)}</i>", style_body))

            # ── P4.2 TOC bookmarks ──
            if columns > 1:
                doc.build(story, onFirstPage=_on_page, onLaterPages=_on_page)
            else:
                # afterFlowable callback for TOC bookmarks
                _bookmark_idx = [0]
                _orig_afterFlowable = getattr(doc, 'afterFlowable', None)
                def _after_flowable(flowable):
                    if _orig_afterFlowable:
                        _orig_afterFlowable(flowable)
                    if toc and hasattr(flowable, 'style'):
                        sname = getattr(flowable.style, 'name', '')
                        if sname in ('LH1', 'LH2', 'LH3'):
                            level = {'LH1': 0, 'LH2': 1, 'LH3': 2}.get(sname, 0)
                            text_val = flowable.getPlainText()
                            key = f"toc_{_bookmark_idx[0]}"
                            _bookmark_idx[0] += 1
                            doc.canv.bookmarkPage(key)
                            doc.canv.addOutlineEntry(text_val, key, level=level)
                doc.afterFlowable = _after_flowable
                doc.build(story, onFirstPage=_on_page, onLaterPages=_on_page)

            # Cleanup temp chart files
            for tp in _chart_temps:
                try:
                    os.unlink(tp)
                except OSError:
                    pass

            logger.info(f"📄 PDF créé: {out_path}")
            return {"success": True, "path": str(out_path), "filename": filename}

        except Exception as e:
            logger.error(f"Erreur create_pdf: {e}")
            return {"success": False, "error": str(e)}

    def create_invoice_pdf(
        self,
        filename: str,
        issuer: Union[Dict[str, str], str],
        client: Union[Dict[str, str], str],
        items: Union[List[Dict], str],
        invoice_meta: Union[Dict[str, str], str, None] = None,
        accent_color: str = "#1a1a2e",
        currency: str = "€",
        page_size: str = "A4",
        orientation: str = "portrait",
        logo_path: Optional[str] = None,
        watermark: Optional[str] = None,
        document_type: str = "facture",
    ) -> Dict[str, Any]:
        """
        Crée une facture PDF professionnelle avec tableau, totaux HT/TVA/TTC.

        Args:
            filename: nom du fichier (ex: "facture-001.pdf")
            issuer: prestataire — dict avec clés: name, address, phone, email, siret (optionnel)
            client: client — dict avec clés: name, address, city (optionnel)
            items: lignes de facture — liste de dicts:
                   {"description": str, "qty": float, "unit_price": float, "vat_rate": float}
                   vat_rate en % (ex: 20 pour 20%). Optionnel: "unit" (libellé unité)
            invoice_meta: dict avec: number, date, due_date, payment_terms, notes (tous optionnels)
            accent_color: couleur principale en hex (ex: "#FF7900" pour orange, "#2d3561" bleu)
            currency: symbole monétaire (ex: "€", "$", "CHF")
        """
        try:
            from reportlab.lib import colors
            from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
            from reportlab.lib.units import cm
            from reportlab.platypus import (
                HRFlowable, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
            )

            issuer = self._parse_json_arg(issuer, {})
            client = self._parse_json_arg(client, {})
            items = self._parse_json_arg(items, [])
            invoice_meta = self._parse_json_arg(invoice_meta, {})
            if not isinstance(invoice_meta, dict):
                invoice_meta = {}

            filename = self._ensure_ext(filename, ".pdf")
            out_path = self._output_path(filename)

            ACCENT = colors.HexColor(accent_color)
            ACCENT_LIGHT = colors.HexColor(
                "#" + "".join(
                    f"{min(255, int(accent_color.lstrip('#')[i:i+2], 16) + 60):02x}"
                    for i in (0, 2, 4)
                )
            )
            WHITE = colors.white
            GREY_BG = colors.HexColor("#f5f5f5")
            GREY_TEXT = colors.HexColor("#555555")
            BLACK = colors.HexColor("#1a1a1a")

            ps = _get_page_size(page_size, orientation)
            doc = SimpleDocTemplate(
                str(out_path), pagesize=ps,
                rightMargin=1.8 * cm, leftMargin=1.8 * cm,
                topMargin=1.5 * cm, bottomMargin=1.5 * cm,
            )
            styles = getSampleStyleSheet()

            def style(name, **kw):
                return ParagraphStyle(name, parent=styles["Normal"], **kw)

            s_company = style("IC", fontSize=18, textColor=ACCENT, leading=22, spaceAfter=2)
            s_issuer  = style("II", fontSize=9,  textColor=GREY_TEXT, leading=13)
            s_label   = style("IL", fontSize=8,  textColor=GREY_TEXT, leading=11)
            s_value   = style("IV", fontSize=10, textColor=BLACK, leading=13)
            s_title   = style("IT", fontSize=28, textColor=ACCENT, alignment=TA_RIGHT, leading=32)
            s_num     = style("IN", fontSize=9,  textColor=GREY_TEXT, alignment=TA_RIGHT, leading=12)
            s_section = style("IS", fontSize=9,  textColor=GREY_TEXT, leading=11, spaceBefore=2)
            s_total_l = style("TL", fontSize=10, textColor=GREY_TEXT, alignment=TA_RIGHT)
            s_total_v = style("TV", fontSize=10, textColor=BLACK, alignment=TA_RIGHT)
            s_grand_l = style("GL", fontSize=12, textColor=WHITE, alignment=TA_RIGHT)
            s_grand_v = style("GV", fontSize=12, textColor=WHITE, alignment=TA_RIGHT)
            s_notes   = style("NO", fontSize=9,  textColor=GREY_TEXT, leading=13)
            s_th      = style("TH", fontSize=9,  textColor=WHITE, leading=12)
            s_td      = style("TD", fontSize=9,  textColor=BLACK, leading=13)
            s_td_r    = style("TDR", fontSize=9, textColor=BLACK, leading=13, alignment=TA_RIGHT)

            story = []
            W = ps[0] - 3.6 * cm  # usable width

            # ── EN-TÊTE : prestataire gauche | FACTURE droite ──
            issuer_block = []

            # Logo intégré dans le bloc prestataire (en haut à gauche)
            if logo_path and Path(logo_path).exists():
                from reportlab.platypus import Image as RLImage
                try:
                    from PIL import Image as PILImage
                    with PILImage.open(logo_path) as _img:
                        _iw, _ih = _img.size
                    # Max 4cm de large, proportionnel
                    _logo_w = min(4 * cm, _iw * 0.75)
                    _logo_h = _logo_w * (_ih / _iw) if _iw else 2 * cm
                    _logo_h = min(_logo_h, 2.5 * cm)
                except Exception:
                    _logo_w, _logo_h = 4 * cm, 2 * cm
                issuer_block.append(RLImage(logo_path, width=_logo_w, height=_logo_h))
                issuer_block.append(Spacer(1, 6))

            issuer_block.append(Paragraph(issuer.get("name", ""), s_company))
            for field in ("address", "city", "phone", "email"):
                val = issuer.get(field, "")
                if val:
                    issuer_block.append(Paragraph(val, s_issuer))
            if issuer.get("siret"):
                issuer_block.append(Paragraph(f"SIRET : {issuer['siret']}", s_issuer))

            num = invoice_meta.get("number", "")
            date_str = invoice_meta.get("date", "")
            due_str  = invoice_meta.get("due_date", "")

            _doc_titles = {
                "facture": "FACTURE", "devis": "DEVIS",
                "bon_commande": "BON DE COMMANDE", "avoir": "AVOIR",
                "note_frais": "NOTE DE FRAIS", "proforma": "FACTURE PROFORMA",
            }
            _doc_title = _doc_titles.get(document_type.lower().strip(), document_type.upper())
            invoice_block = [
                Paragraph(_doc_title, s_title),
            ]
            if num:
                invoice_block.append(Paragraph(f"N° {num}", s_num))
            if date_str:
                invoice_block.append(Paragraph(f"Date : {date_str}", s_num))
            if due_str:
                invoice_block.append(Paragraph(f"Échéance : {due_str}", s_num))

            header_table = Table(
                [[issuer_block, invoice_block]],
                colWidths=[W * 0.55, W * 0.45],
            )
            header_table.setStyle(TableStyle([
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING",  (0, 0), (-1, -1), 0),
                ("RIGHTPADDING", (0, 0), (-1, -1), 0),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 0),
                ("TOPPADDING",   (0, 0), (-1, -1), 0),
            ]))
            story.append(header_table)
            story.append(HRFlowable(width="100%", thickness=1.5, color=ACCENT, spaceBefore=10, spaceAfter=14))

            # ── BLOC CLIENT ──
            _client_label = "FACTURER À" if document_type.lower() == "facture" else "CLIENT"
            client_lines = [Paragraph(_client_label, s_label)]
            if client.get("name"):
                client_lines.append(Paragraph(f"<b>{client['name']}</b>", s_value))
            for f in ("address", "city", "email", "phone"):
                v = client.get(f, "")
                if v:
                    client_lines.append(Paragraph(v, s_section))

            client_table = Table([[client_lines]], colWidths=[W])
            client_table.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (0, 0), GREY_BG),
                ("LEFTPADDING",  (0, 0), (0, 0), 10),
                ("RIGHTPADDING", (0, 0), (0, 0), 10),
                ("TOPPADDING",   (0, 0), (0, 0), 8),
                ("BOTTOMPADDING",(0, 0), (0, 0), 8),
                ("ROUNDEDCORNERS", [4]),
            ]))
            story.append(client_table)
            story.append(Spacer(1, 16))

            # ── TABLEAU DES LIGNES ──
            col_w = [W * 0.44, W * 0.10, W * 0.16, W * 0.10, W * 0.10, W * 0.10]
            th = [
                Paragraph("DÉSIGNATION", s_th),
                Paragraph("QTÉ", s_th),
                Paragraph(f"P.U. HT ({currency})", s_th),
                Paragraph("TVA", s_th),
                Paragraph(f"TVA ({currency})", s_th),
                Paragraph(f"TOTAL TTC ({currency})", s_th),
            ]
            table_data = [th]

            subtotal_ht = 0.0
            vat_buckets: Dict[float, float] = {}

            for i, item in enumerate(items):
                if isinstance(item, str):
                    continue
                desc       = item.get("description") or item.get("designation") or item.get("nom") or item.get("label") or ""
                qty        = float(item.get("qty") or item.get("quantity") or item.get("quantite") or 1)
                unit_price = float(item.get("unit_price") or item.get("price") or item.get("prix") or item.get("prix_unitaire") or 0)
                vat_rate   = float(item.get("vat_rate") or item.get("tva") or item.get("tax") or 20)
                unit_label = item.get("unit") or item.get("unite") or ""

                ht      = qty * unit_price
                vat_amt = ht * vat_rate / 100
                ttc     = ht + vat_amt

                subtotal_ht += ht
                vat_buckets[vat_rate] = vat_buckets.get(vat_rate, 0.0) + vat_amt

                qty_str = f"{qty:g}" + (f" {unit_label}" if unit_label else "")
                bg = WHITE if i % 2 == 0 else GREY_BG

                row = [
                    Paragraph(desc, s_td),
                    Paragraph(qty_str, s_td_r),
                    Paragraph(f"{unit_price:,.2f}", s_td_r),
                    Paragraph(f"{vat_rate:g}%", s_td_r),
                    Paragraph(f"{vat_amt:,.2f}", s_td_r),
                    Paragraph(f"{ttc:,.2f}", s_td_r),
                ]
                table_data.append(row)

            items_table = Table(table_data, colWidths=col_w, repeatRows=1)
            ts = TableStyle([
                # En-tête
                ("BACKGROUND",    (0, 0), (-1, 0), ACCENT),
                ("FONTNAME",      (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE",      (0, 0), (-1, 0), 9),
                ("ROWBACKGROUNDS",(0, 1), (-1, -1), [WHITE, GREY_BG]),
                ("FONTSIZE",      (0, 1), (-1, -1), 9),
                ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
                ("LEFTPADDING",   (0, 0), (-1, -1), 6),
                ("RIGHTPADDING",  (0, 0), (-1, -1), 6),
                ("TOPPADDING",    (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                ("LINEBELOW",     (0, 0), (-1, 0), 0, ACCENT),
                ("LINEBELOW",     (0, 1), (-1, -1), 0.3, colors.HexColor("#dddddd")),
                ("BOX",           (0, 0), (-1, -1), 0.5, colors.HexColor("#dddddd")),
            ])
            items_table.setStyle(ts)
            story.append(items_table)
            story.append(Spacer(1, 10))

            # ── TOTAUX ──
            total_vat = sum(vat_buckets.values())
            total_ttc = subtotal_ht + total_vat

            totals_rows = []
            totals_rows.append([Paragraph("Sous-total HT :", s_total_l),
                                 Paragraph(f"{subtotal_ht:,.2f} {currency}", s_total_v)])
            for rate, amt in sorted(vat_buckets.items()):
                totals_rows.append([
                    Paragraph(f"TVA {rate:g}% :", s_total_l),
                    Paragraph(f"{amt:,.2f} {currency}", s_total_v),
                ])
            totals_rows.append([Paragraph("TOTAL TTC :", s_grand_l),
                                 Paragraph(f"{total_ttc:,.2f} {currency}", s_grand_v)])

            totals_table = Table(totals_rows, colWidths=[W * 0.75, W * 0.25])
            n = len(totals_rows)
            totals_table.setStyle(TableStyle([
                ("ALIGN",         (0, 0), (-1, -1), "RIGHT"),
                ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
                ("TOPPADDING",    (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ("LEFTPADDING",   (0, 0), (-1, -1), 6),
                ("RIGHTPADDING",  (0, 0), (-1, -1), 6),
                ("BACKGROUND",    (0, n-1), (-1, n-1), ACCENT),
                ("LINEABOVE",     (0, n-1), (-1, n-1), 1, ACCENT),
            ]))
            story.append(totals_table)

            # ── NOTES / CONDITIONS ──
            notes = invoice_meta.get("notes", "")
            payment = invoice_meta.get("payment_terms", "")
            if notes or payment:
                story.append(Spacer(1, 14))
                story.append(HRFlowable(width="100%", thickness=0.5,
                                        color=colors.HexColor("#cccccc"), spaceAfter=8))
                if payment:
                    story.append(Paragraph(f"<b>Conditions de paiement :</b> {payment}", s_notes))
                if notes:
                    story.append(Paragraph(f"<b>Notes :</b> {notes}", s_notes))

            # ── Watermark (optionnel) ──
            def _invoice_on_page(canvas_obj, doc_obj):
                if watermark:
                    canvas_obj.saveState()
                    canvas_obj.setFont("Helvetica-Bold", 48)
                    canvas_obj.setFillGray(0.85)
                    canvas_obj.translate(ps[0] / 2, ps[1] / 2)
                    canvas_obj.rotate(45)
                    canvas_obj.drawCentredString(0, 0, watermark)
                    canvas_obj.restoreState()

            doc.build(story, onFirstPage=_invoice_on_page, onLaterPages=_invoice_on_page)
            logger.info(f"📄 Facture PDF créée: {out_path}")
            return {
                "success": True,
                "path": str(out_path),
                "filename": filename,
                "total_ttc": f"{total_ttc:,.2f} {currency}",
                "invoice_meta": invoice_meta,
            }

        except Exception as e:
            logger.error(f"Erreur create_invoice_pdf: {e}")
            return {"success": False, "error": str(e)}

    # ──────────────────────────────────────────────────────
    # DOCX (via python-docx)
    # ──────────────────────────────────────────────────────

    @staticmethod
    def _docx_add_runs(paragraph, text: str) -> None:
        """Ajoute des runs avec support markdown inline."""
        parts = re.split(r"(\*\*\*.*?\*\*\*|\*\*.*?\*\*|\*[^*\n]+?\*|`[^`]+?`)", text)
        for part in parts:
            if part.startswith("***") and part.endswith("***") and len(part) > 6:
                run = paragraph.add_run(part[3:-3])
                run.bold = True
                run.italic = True
            elif part.startswith("**") and part.endswith("**") and len(part) > 4:
                run = paragraph.add_run(part[2:-2])
                run.bold = True
            elif part.startswith("*") and part.endswith("*") and len(part) > 2:
                run = paragraph.add_run(part[1:-1])
                run.italic = True
            elif part.startswith("`") and part.endswith("`") and len(part) > 2:
                from docx.shared import Pt
                run = paragraph.add_run(part[1:-1])
                run.font.name = "Courier New"
                run.font.size = Pt(10)
            else:
                paragraph.add_run(part)

    @staticmethod
    def _docx_add_hr(paragraph) -> None:
        """Ajoute une ligne horizontale (bordure basse) à un paragraphe."""
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn

        pPr = paragraph._p.get_or_add_pPr()
        pBdr = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), "6")
        bottom.set(qn("w:space"), "1")
        bottom.set(qn("w:color"), "AAAAAA")
        pBdr.append(bottom)
        pPr.append(pBdr)

    def create_docx(
        self,
        filename: str,
        title: str,
        content: str,
        author: str = "Lumena",
        images: Optional[List[Dict]] = None,
        header_footer: Optional[Dict] = None,
        page_size: str = "A4",
        orientation: str = "portrait",
        toc: bool = False,
        charts: Optional[List[Dict]] = None,
        theme: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Crée un fichier Word .docx depuis du contenu texte/markdown.
        Supporte : # ## ### titres, **gras**, *italique*, `code`, listes -, 1.,
        séparateurs ---, tables, checkboxes, footnotes, TOC, thèmes, charts inline.
        """
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.oxml import OxmlElement
            from docx.oxml.ns import qn
            from docx.shared import Cm, Mm, Pt, RGBColor

            images = self._parse_json_arg(images, [])
            header_footer = self._parse_json_arg(header_footer, None)
            charts = self._parse_json_arg(charts, [])

            filename = self._ensure_ext(filename, ".docx")
            out_path = self._output_path(filename)

            # ── P4.1 Theme ──
            th = self._get_theme(theme)

            def _hex_to_rgb(h: str) -> RGBColor:
                h = h.lstrip("#")
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

            # ── P4.5 Footnotes ──
            content, footnotes_dict = self._parse_footnotes(content)
            _fn_counter = [0]

            def _add_footnote_ref(paragraph, key: str):
                """Ajoute une référence de footnote OOXML."""
                if key not in footnotes_dict:
                    return
                _fn_counter[0] += 1
                fn_id = _fn_counter[0]
                # Superscript number in text
                run = paragraph.add_run(f"{fn_id}")
                run.font.superscript = True
                run.font.size = Pt(8)

            doc = Document()

            # Page size / orientation
            _PAGE_DIMS = {
                "A4": (Mm(210), Mm(297)), "A3": (Mm(297), Mm(420)),
                "A5": (Mm(148), Mm(210)), "LETTER": (Mm(216), Mm(279)),
                "LEGAL": (Mm(216), Mm(356)),
            }
            pw, ph = _PAGE_DIMS.get(page_size.upper(), _PAGE_DIMS["A4"])
            if orientation.lower() == "landscape":
                pw, ph = ph, pw
            for section in doc.sections:
                section.page_width = pw
                section.page_height = ph
                section.top_margin = Cm(2.5)
                section.bottom_margin = Cm(2)
                section.left_margin = Cm(2.5)
                section.right_margin = Cm(2.5)

            # Header/footer
            if header_footer:
                section = doc.sections[0]
                hl = header_footer.get("header_left", "")
                hr_ = header_footer.get("header_right", "")
                if hl or hr_:
                    header = section.header
                    header.is_linked_to_previous = False
                    hp = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
                    hp.text = hl
                    if hr_:
                        hp.text += "\t\t" + hr_
                fc = header_footer.get("footer_center", "")
                if fc:
                    footer = section.footer
                    footer.is_linked_to_previous = False
                    fp = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
                    fp.text = fc.replace("{page}", "").replace("{total}", "")
                    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Métadonnées
            doc.core_properties.author = author
            doc.core_properties.title = title

            # ── P4.2 TOC field ──
            if toc:
                toc_para = doc.add_paragraph()
                toc_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
                run_toc = toc_para.add_run("Table des matières")
                run_toc.bold = True
                run_toc.font.size = Pt(th["font_size_h2"])
                run_toc.font.color.rgb = _hex_to_rgb(th["primary"])
                # Insert TOC field code
                fld_begin = OxmlElement("w:fldSimple")
                fld_begin.set(qn("w:instr"), r'TOC \o "1-3" \h \z \u')
                toc_para._p.append(fld_begin)
                doc.add_paragraph("")  # spacer

            # Titre principal
            title_para = doc.add_heading(title, level=0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            # Apply theme color to title
            for run in title_para.runs:
                run.font.color.rgb = _hex_to_rgb(th["primary"])

            # Parse tables
            all_lines, md_tables = _parse_md_tables(content)
            table_line_ranges = set()
            for (ts, te, _, _) in md_tables:
                table_line_ranges.update(range(ts, te))

            def _process_footnote_refs_in_para(p, text: str):
                """Parse [^N] refs in text and add runs + footnote refs."""
                parts = re.split(r"(\[\^\w+\])", text)
                for part in parts:
                    m = re.match(r"^\[\^(\w+)\]$", part)
                    if m:
                        _add_footnote_ref(p, m.group(1))
                    elif part:
                        self._docx_add_runs(p, part)

            for line_idx, line in enumerate(all_lines):
                if line_idx in table_line_ranges:
                    for (ts, te, t_headers, t_rows) in md_tables:
                        if line_idx == ts:
                            n_cols = len(t_headers)
                            n_rows = len(t_rows) + 1
                            tbl = doc.add_table(rows=n_rows, cols=n_cols, style="Table Grid")
                            for ci, h in enumerate(t_headers):
                                cell = tbl.rows[0].cells[ci]
                                cell.text = h
                                for run in cell.paragraphs[0].runs:
                                    run.bold = True
                            for ri, row in enumerate(t_rows):
                                for ci, val in enumerate(row):
                                    if ci < n_cols:
                                        tbl.rows[ri + 1].cells[ci].text = val
                    continue

                s = line.strip()
                if not s:
                    doc.add_paragraph("")
                    continue

                # Inline image
                img_match = re.match(r"^!\[([^\]]*)\]\(([^)]+)\)\s*$", s)
                if img_match:
                    caption, img_path = img_match.group(1), img_match.group(2)
                    if Path(img_path).exists():
                        doc.add_picture(img_path, width=Cm(12))
                        if caption:
                            cp = doc.add_paragraph()
                            cr = cp.add_run(caption)
                            cr.italic = True
                    continue

                # ── P4.4 Checkboxes ──
                cb_match = re.match(r"^- \[([ xX])\]\s+(.*)", s)
                if cb_match:
                    checked = cb_match.group(1).lower() == "x"
                    cb_char = "☑" if checked else "☐"
                    p = doc.add_paragraph(style="List Bullet")
                    p.text = ""
                    r = p.add_run(f"{cb_char} ")
                    self._docx_add_runs(p, cb_match.group(2))
                    continue

                if s.startswith("### "):
                    h = doc.add_heading(s[4:], level=3)
                    for run in h.runs:
                        run.font.color.rgb = _hex_to_rgb(th.get("text", "#4a5568"))
                elif s.startswith("## "):
                    h = doc.add_heading(s[3:], level=2)
                    for run in h.runs:
                        run.font.color.rgb = _hex_to_rgb(th["secondary"])
                elif s.startswith("# "):
                    h = doc.add_heading(s[2:], level=1)
                    for run in h.runs:
                        run.font.color.rgb = _hex_to_rgb(th["primary"])
                elif s.startswith(("- ", "* ", "• ")):
                    p = doc.add_paragraph(style="List Bullet")
                    _process_footnote_refs_in_para(p, s[2:])
                elif re.match(r"^\d+\.\s", s):
                    p = doc.add_paragraph(style="List Number")
                    _process_footnote_refs_in_para(p, re.sub(r"^\d+\.\s+", "", s))
                elif re.match(r"^-{3,}$|^={3,}$", s):
                    p = doc.add_paragraph()
                    self._docx_add_hr(p)
                else:
                    p = doc.add_paragraph()
                    _process_footnote_refs_in_para(p, s)

            # Append explicit images
            for img_def in (images or []):
                img_path = img_def.get("path", "")
                if img_path and Path(img_path).exists():
                    w_cm = img_def.get("width_cm", 12)
                    doc.add_picture(img_path, width=Cm(w_cm))
                    cap = img_def.get("caption", "")
                    if cap:
                        cp = doc.add_paragraph()
                        cr = cp.add_run(cap)
                        cr.italic = True

            # ── P4.3 Inline charts ──
            _chart_temps = []
            for chart_def in (charts or []):
                chart_path = chart_def.get("path", "")
                if not chart_path or not Path(chart_path).exists():
                    gen_path = self._make_inline_chart(chart_def)
                    if gen_path:
                        chart_path = gen_path
                        _chart_temps.append(gen_path)
                if chart_path and Path(chart_path).exists():
                    w_cm = chart_def.get("width_cm", 14)
                    doc.add_picture(chart_path, width=Cm(w_cm))
                    cap = chart_def.get("caption", chart_def.get("title", ""))
                    if cap:
                        cp = doc.add_paragraph()
                        cr = cp.add_run(cap)
                        cr.italic = True

            # ── P4.5 Footnotes endnotes section ──
            if footnotes_dict:
                doc.add_paragraph("")
                hr_p = doc.add_paragraph()
                self._docx_add_hr(hr_p)
                fn_title = doc.add_paragraph()
                r = fn_title.add_run("Notes")
                r.bold = True
                r.font.size = Pt(10)
                for i, (key, text) in enumerate(footnotes_dict.items(), 1):
                    fn_p = doc.add_paragraph()
                    rn = fn_p.add_run(f"{i}. ")
                    rn.font.size = Pt(9)
                    rn.font.superscript = True
                    rt = fn_p.add_run(text)
                    rt.font.size = Pt(9)

            doc.save(str(out_path))

            # Cleanup temp chart files
            for tp in _chart_temps:
                try:
                    os.unlink(tp)
                except OSError:
                    pass

            logger.info(f"📝 DOCX créé: {out_path}")
            return {"success": True, "path": str(out_path), "filename": filename}

        except Exception as e:
            logger.error(f"Erreur create_docx: {e}")
            return {"success": False, "error": str(e)}

    # ──────────────────────────────────────────────────────
    # XLSX (via openpyxl)
    # ──────────────────────────────────────────────────────

    def create_xlsx(
        self,
        filename: str,
        sheets: Any,
        author: str = "Lumena",
    ) -> Dict[str, Any]:
        """
        Crée un fichier Excel .xlsx.

        sheets: liste (ou JSON string) de feuilles au format :
          [{"name": "NomFeuille", "title": "opt", "headers": ["Col1","Col2"], "rows": [[v1,v2], ...]}]

        Exemple minimal :
          [{"name": "Stats", "headers": ["Nom","Valeur"], "rows": [["Total",42]]}]
        """
        try:
            import openpyxl
            from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
            from openpyxl.utils import get_column_letter

            filename = self._ensure_ext(filename, ".xlsx")
            out_path = self._output_path(filename)

            sheets_data = self._parse_json_arg(sheets, fallback=[])
            if isinstance(sheets_data, dict):
                sheets_data = [sheets_data]
            if not sheets_data:
                return {"success": False, "error": "Aucune donnée de feuille fournie."}

            wb = openpyxl.Workbook()
            wb.properties.creator = author
            wb.remove(wb.active)

            # Styles
            HEADER_FILL = PatternFill("solid", fgColor="1A1A2E")
            HEADER_FONT = Font(bold=True, color="FFFFFF", size=11, name="Arial")
            TITLE_FONT  = Font(bold=True, size=13, color="1A1A2E", name="Arial")
            DATA_FONT   = Font(size=10, name="Arial")
            BORDER_SIDE = Side(style="thin", color="CCCCCC")
            CELL_BORDER = Border(
                left=BORDER_SIDE, right=BORDER_SIDE,
                top=BORDER_SIDE,  bottom=BORDER_SIDE,
            )
            ALT_FILL = PatternFill("solid", fgColor="F0F4FF")

            for sheet_def in sheets_data:
                name     = str(sheet_def.get("name", "Feuille"))[:31]
                headers  = sheet_def.get("headers", [])
                rows     = sheet_def.get("rows", [])
                s_title  = sheet_def.get("title", "")

                ws = wb.create_sheet(title=name)
                row_offset = 1

                if s_title:
                    nb_cols = max(len(headers), 1)
                    cell = ws.cell(row=1, column=1, value=s_title)
                    cell.font = TITLE_FONT
                    cell.alignment = Alignment(horizontal="center", vertical="center")
                    if nb_cols > 1:
                        ws.merge_cells(
                            start_row=1, start_column=1,
                            end_row=1,   end_column=nb_cols,
                        )
                    ws.row_dimensions[1].height = 22
                    row_offset = 2

                if headers:
                    for col_idx, header in enumerate(headers, start=1):
                        cell = ws.cell(row=row_offset, column=col_idx, value=header)
                        cell.font = HEADER_FONT
                        cell.fill = HEADER_FILL
                        cell.alignment = Alignment(horizontal="center", vertical="center")
                        cell.border = CELL_BORDER
                    ws.row_dimensions[row_offset].height = 20
                    row_offset += 1

                for r_idx, row in enumerate(rows):
                    for c_idx, value in enumerate(row, start=1):
                        cell = ws.cell(row=row_offset + r_idx, column=c_idx, value=value)
                        cell.font = DATA_FONT
                        cell.border = CELL_BORDER
                        cell.alignment = Alignment(vertical="center")
                        if r_idx % 2 == 1:
                            cell.fill = ALT_FILL

                # Auto-largeur des colonnes
                nb_cols = len(headers) if headers else (max((len(r) for r in rows), default=1))
                for col_idx in range(1, nb_cols + 1):
                    col_letter = get_column_letter(col_idx)
                    max_len = 0
                    for row_cells in ws.iter_rows(min_col=col_idx, max_col=col_idx):
                        for c in row_cells:
                            if c.value is not None:
                                max_len = max(max_len, len(str(c.value)))
                    ws.column_dimensions[col_letter].width = min(max_len + 4, 50)

            wb.save(str(out_path))
            logger.info(f"📊 XLSX créé: {out_path}")
            return {"success": True, "path": str(out_path), "filename": filename}

        except Exception as e:
            logger.error(f"Erreur create_xlsx: {e}")
            return {"success": False, "error": str(e)}

    # ──────────────────────────────────────────────────────
    # PPTX (via python-pptx)
    # ──────────────────────────────────────────────────────

    def create_pptx(
        self,
        filename: str,
        title: str,
        slides: Any,
        theme_color: str = "1a1a2e",
        images: Optional[List[Dict[str, str]]] = None,
    ) -> Dict[str, Any]:
        """
        Crée une présentation .pptx.

        slides: liste (ou JSON string) de slides au format :
          [{"title": "Titre slide", "content": "texte markdown"}]

        Supporte dans content : ## sous-titre, - bullet, texte libre.
        theme_color: couleur hex sans # pour le fond titre/header (défaut: 1a1a2e = bleu nuit).
        """
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Cm, Inches, Pt

            filename = self._ensure_ext(filename, ".pptx")
            out_path = self._output_path(filename)

            slides_data = self._parse_json_arg(slides, fallback=[])
            if isinstance(slides_data, dict):
                slides_data = [slides_data]
            if not slides_data:
                return {"success": False, "error": "Aucune donnée de slide fournie."}

            # Couleurs — robuste : accepte hex (#1a1a2e / 1a1a2e) ou noms CSS
            import re as _re
            _CSS_COLORS = {
                "black": "000000", "white": "ffffff", "red": "ff0000",
                "green": "008000", "blue": "0000ff", "navy": "000080",
                "darkblue": "00008b", "midnightblue": "191970",
                "cobaltblue": "0047ab", "royalblue": "4169e1",
                "steelblue": "4682b4", "dodgerblue": "1e90ff",
                "darkred": "8b0000", "crimson": "dc143c",
                "maroon": "800000", "darkgreen": "006400",
                "teal": "008080", "purple": "800080",
                "indigo": "4b0082", "gold": "ffd700",
                "orange": "ffa500", "coral": "ff7f50",
                "gray": "808080", "grey": "808080",
                "silver": "c0c0c0", "slategray": "708090",
            }
            hex_clean = theme_color.strip().lstrip("#").strip()
            # Si c'est un nom CSS, convertir
            if not _re.fullmatch(r"[0-9a-fA-F]{6}", hex_clean):
                hex_clean = _CSS_COLORS.get(hex_clean.lower(), "1a1a2e")
            primary   = RGBColor(
                int(hex_clean[0:2], 16),
                int(hex_clean[2:4], 16),
                int(hex_clean[4:6], 16),
            )
            white     = RGBColor(0xFF, 0xFF, 0xFF)
            light_bg  = RGBColor(0xF5, 0xF7, 0xFF)
            text_dark = RGBColor(0x1A, 0x1A, 0x2E)
            accent    = RGBColor(0x4A, 0x9E, 0xFF)
            subtitle_color = RGBColor(0xCC, 0xCC, 0xFF)

            prs = Presentation()
            prs.slide_width  = Inches(13.33)
            prs.slide_height = Inches(7.5)
            blank = prs.slide_layouts[6]  # Layout vide

            # ── Slide de titre ──────────────────────────────
            title_slide = prs.slides.add_slide(blank)
            bg = title_slide.background
            bg.fill.solid()
            bg.fill.fore_color.rgb = primary

            # Titre
            tb = title_slide.shapes.add_textbox(
                Inches(1), Inches(2.3), Inches(11.33), Inches(1.8)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = title
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = white

            # Sous-titre (date)
            tb2 = title_slide.shapes.add_textbox(
                Inches(1), Inches(4.3), Inches(11.33), Inches(0.7)
            )
            p2 = tb2.text_frame.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = f"Généré par Lumena • {datetime.now().strftime('%d/%m/%Y')}"
            r2.font.size = Pt(16)
            r2.font.color.rgb = subtitle_color

            # ── Slides de contenu ───────────────────────────
            for slide_def in slides_data:
                s_title   = slide_def.get("title", "")
                s_content = slide_def.get("content", "")

                slide = prs.slides.add_slide(blank)
                bg = slide.background
                bg.fill.solid()
                bg.fill.fore_color.rgb = light_bg

                # Barre header
                header_shape = slide.shapes.add_shape(
                    1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
                    Inches(0), Inches(0),
                    prs.slide_width, Inches(1.1),
                )
                header_shape.fill.solid()
                header_shape.fill.fore_color.rgb = primary
                header_shape.line.fill.background()

                # Titre de la slide
                tb_title = slide.shapes.add_textbox(
                    Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.9)
                )
                p_title = tb_title.text_frame.paragraphs[0]
                r_title = p_title.add_run()
                r_title.text = s_title
                r_title.font.size = Pt(24)
                r_title.font.bold = True
                r_title.font.color.rgb = white

                # Zone de contenu
                tb_content = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.25), Inches(12.3), Inches(5.9)
                )
                tf_content = tb_content.text_frame
                tf_content.word_wrap = True

                first_para = True
                for line in s_content.strip().split("\n"):
                    s = line.strip()
                    if not s:
                        continue

                    if first_para:
                        p = tf_content.paragraphs[0]
                        first_para = False
                    else:
                        p = tf_content.add_paragraph()

                    if s.startswith("## "):
                        r = p.add_run()
                        r.text = s[3:]
                        r.font.size = Pt(20)
                        r.font.bold = True
                        r.font.color.rgb = primary
                        p.space_before = Pt(8)
                    elif s.startswith("### "):
                        r = p.add_run()
                        r.text = s[4:]
                        r.font.size = Pt(16)
                        r.font.bold = True
                        r.font.color.rgb = text_dark
                        p.space_before = Pt(5)
                    elif s.startswith(("- ", "* ", "• ")):
                        p.level = 1
                        r = p.add_run()
                        r.text = "▸  " + s[2:]
                        r.font.size = Pt(14)
                        r.font.color.rgb = text_dark
                        p.space_before = Pt(4)
                    else:
                        r = p.add_run()
                        r.text = s
                        r.font.size = Pt(13)
                        r.font.color.rgb = text_dark
                        p.space_before = Pt(3)

            # ── Images additionnelles (images param) ──
            if images:
                for img_def in (images if isinstance(images, list) else []):
                    img_path = img_def.get("path", "")
                    caption = img_def.get("caption", "")
                    if not img_path or not Path(img_path).exists():
                        continue
                    slide = prs.slides.add_slide(blank)
                    bg = slide.background
                    bg.fill.solid()
                    bg.fill.fore_color.rgb = light_bg
                    # Image centrée
                    slide.shapes.add_picture(
                        img_path, Inches(1.5), Inches(0.8), Inches(10), Inches(5.5)
                    )
                    if caption:
                        tb_cap = slide.shapes.add_textbox(
                            Inches(1.5), Inches(6.5), Inches(10), Inches(0.7)
                        )
                        p_cap = tb_cap.text_frame.paragraphs[0]
                        p_cap.alignment = PP_ALIGN.CENTER
                        r_cap = p_cap.add_run()
                        r_cap.text = caption
                        r_cap.font.size = Pt(12)
                        r_cap.font.color.rgb = text_dark

            prs.save(str(out_path))
            logger.info(f"🎞️ PPTX créé: {out_path}")
            return {"success": True, "path": str(out_path), "filename": filename}

        except Exception as e:
            logger.error(f"Erreur create_pptx: {e}")
            return {"success": False, "error": str(e)}

    # ──────────────────────────────────────────────────────
    # Templates Jinja2 — P2
    # ──────────────────────────────────────────────────────

    @property
    def _TEMPLATES_DIR(self) -> Path:
        from src.utils.paths import TEMPLATES_DIR
        return TEMPLATES_DIR

    def list_templates(self) -> Dict[str, Any]:
        """Liste tous les templates disponibles (builtin + custom) avec leurs variables."""
        try:
            import re as _re
            templates = []
            tdir = self._TEMPLATES_DIR
            if tdir.exists():
                for f in sorted(tdir.glob("*.html.j2")):
                    variables = self._extract_template_variables(f)
                    templates.append({"name": f.stem.replace(".html", ""), "path": str(f), "type": "builtin", "variables": variables})
                custom_dir = tdir / "custom"
                if custom_dir.exists():
                    for f in sorted(custom_dir.glob("*.html.j2")):
                        variables = self._extract_template_variables(f)
                        templates.append({"name": f.stem.replace(".html", ""), "path": str(f), "type": "custom", "variables": variables})
            return {"success": True, "templates": templates, "count": len(templates)}
        except Exception as e:
            return {"success": False, "error": str(e)}

    @staticmethod
    def _extract_template_variables(template_path: Path) -> List[str]:
        """Extrait les noms de variables Jinja2 d'un template (top-level uniquement)."""
        import re as _re
        try:
            content = template_path.read_text(encoding="utf-8")
            # Capture {{ var }}, {{ var.attr }}, {{ var|filter }}
            raw = _re.findall(r'\{\{\s*([a-zA-Z_][a-zA-Z0-9_]*)(?:\.[a-zA-Z_][a-zA-Z0-9_]*|\|[^}]*)?\s*\}\}', content)
            # Also capture {% for x in var %}
            raw += _re.findall(r'\{%[^%]*\bfor\s+\w+\s+in\s+([a-zA-Z_][a-zA-Z0-9_]*)\b', content)
            # Also capture {% if var %}
            raw += _re.findall(r'\{%\s*if\s+([a-zA-Z_][a-zA-Z0-9_]*)\b', content)
            # Exclude Jinja2 builtins, loop vars, and single-char vars (loop artifacts)
            _BUILTINS = {"loop", "true", "false", "none", "range", "lipsum", "dict", "cycler", "joiner", "namespace", "not", "item"}
            seen = set()
            result = []
            for v in raw:
                vl = v.lower()
                if len(v) <= 1 or vl in _BUILTINS or v not in v:
                    continue
                if v not in seen:
                    seen.add(v)
                    result.append(v)
            return sorted(result)
        except Exception:
            return []

    def create_from_template(
        self,
        template_name: str,
        variables: Union[Dict, str],
        output_filename: str,
        output_format: str = "pdf",
    ) -> Dict[str, Any]:
        """Rend un template Jinja2 avec les variables fournies et produit PDF ou HTML."""
        try:
            from jinja2 import Environment, FileSystemLoader

            variables = self._parse_json_arg(variables, {})
            tdir = self._TEMPLATES_DIR
            # Chercher dans builtin puis custom
            tfile = tdir / f"{template_name}.html.j2"
            if not tfile.exists():
                tfile = tdir / "custom" / f"{template_name}.html.j2"
            if not tfile.exists():
                return {"success": False, "error": f"Template '{template_name}' non trouvé."}

            env = Environment(
                loader=FileSystemLoader([str(tdir), str(tdir / "custom")]),
                autoescape=True,
            )
            template = env.get_template(tfile.name)
            html = template.render(**variables)

            fmt = output_format.lower().lstrip(".")
            if fmt == "html":
                output_filename = self._ensure_ext(output_filename, ".html")
                out_path = self._output_path(output_filename)
                out_path.write_text(html, encoding="utf-8")
            else:
                output_filename = self._ensure_ext(output_filename, ".pdf")
                out_path = self._output_path(output_filename)
                # Try WeasyPrint first, fallback to Playwright
                try:
                    import weasyprint
                    weasyprint.HTML(string=html).write_pdf(str(out_path))
                except Exception as wp_err:
                    logger.warning(f"WeasyPrint indisponible ({wp_err}), fallback Playwright")
                    from playwright.sync_api import sync_playwright
                    with sync_playwright() as p:
                        browser = p.chromium.launch(headless=True)
                        page = browser.new_page()
                        page.set_content(html, wait_until="networkidle")
                        page.pdf(path=str(out_path), format="A4", print_background=True)
                        browser.close()

            logger.info(f"📄 Template '{template_name}' → {out_path}")
            return {"success": True, "path": str(out_path), "filename": output_filename, "template": template_name}
        except Exception as e:
            logger.error(f"Erreur create_from_template: {e}")
            return {"success": False, "error": str(e)}

    def save_template(
        self, template_name: str, html_content: str
    ) -> Dict[str, Any]:
        """Sauvegarde un template HTML Jinja2 custom."""
        try:
            custom_dir = self._TEMPLATES_DIR / "custom"
            custom_dir.mkdir(parents=True, exist_ok=True)
            safe_name = re.sub(r"[^\w\-]", "_", template_name)
            tfile = custom_dir / f"{safe_name}.html.j2"
            tfile.write_text(html_content, encoding="utf-8")
            logger.info(f"📝 Template custom sauvé: {tfile}")
            return {"success": True, "path": str(tfile), "name": safe_name}
        except Exception as e:
            logger.error(f"Erreur save_template: {e}")
            return {"success": False, "error": str(e)}

    # ──────────────────────────────────────────────────────
    # P3 — Intelligence documentaire avancée
    # ──────────────────────────────────────────────────────

    def add_watermark(
        self, input_path: str, text: str = "CONFIDENTIEL",
        opacity: float = 0.15, angle: float = 45,
        output_path: Optional[str] = None,
    ) -> Dict[str, Any]:
        """Ajoute un filigrane texte sur toutes les pages d'un PDF."""
        try:
            from pypdf import PdfReader, PdfWriter
            from reportlab.lib.pagesizes import A4
            from reportlab.pdfgen.canvas import Canvas

            pp = Path(input_path)
            if not pp.exists():
                found = self._find_in_workspace(input_path)
                if found:
                    pp = found
                else:
                    return {"success": False, "error": f"Fichier non trouvé: {input_path}"}

            reader = PdfReader(str(pp))
            writer = PdfWriter()

            for page in reader.pages:
                mbox = page.mediabox
                pw, ph = float(mbox.width), float(mbox.height)
                buf = io.BytesIO()
                c = Canvas(buf, pagesize=(pw, ph))
                c.saveState()
                c.setFont("Helvetica-Bold", 60)
                c.setFillGray(0.5, alpha=opacity)
                c.translate(pw / 2, ph / 2)
                c.rotate(angle)
                c.drawCentredString(0, 0, text)
                c.restoreState()
                c.save()
                buf.seek(0)
                overlay = PdfReader(buf)
                page.merge_page(overlay.pages[0])
                writer.add_page(page)

            out = Path(output_path) if output_path else self._output_path(f"{pp.stem}_watermarked.pdf")
            out.parent.mkdir(parents=True, exist_ok=True)
            with open(str(out), "wb") as f:
                writer.write(f)
            logger.info(f"📄 Watermark '{text}' ajouté: {out}")
            return {"success": True, "path": str(out), "filename": out.name}
        except Exception as e:
            logger.error(f"Erreur add_watermark: {e}")
            return {"success": False, "error": str(e)}

    def add_signature(
        self, input_path: str, signature_image_path: str,
        page: int = -1, position: str = "bottom-right",
        date_text: str = "auto", output_path: Optional[str] = None,
    ) -> Dict[str, Any]:
        """Ajoute une image de signature sur une page d'un PDF."""
        try:
            from pypdf import PdfReader, PdfWriter
            from reportlab.lib.units import cm
            from reportlab.pdfgen.canvas import Canvas

            pp = Path(input_path)
            if not pp.exists():
                found = self._find_in_workspace(input_path)
                if found:
                    pp = found
                else:
                    return {"success": False, "error": f"Fichier non trouvé: {input_path}"}

            sig_path = Path(signature_image_path)
            if not sig_path.exists():
                return {"success": False, "error": f"Image signature non trouvée: {signature_image_path}"}

            reader = PdfReader(str(pp))
            writer = PdfWriter()
            total = len(reader.pages)
            target_page = page if page >= 0 else total + page

            for idx, pg in enumerate(reader.pages):
                if idx == target_page:
                    mbox = pg.mediabox
                    pw, ph = float(mbox.width), float(mbox.height)
                    buf = io.BytesIO()
                    c = Canvas(buf, pagesize=(pw, ph))

                    sig_w, sig_h = 4 * cm, 2 * cm
                    if position == "bottom-right":
                        x, y = pw - sig_w - 40, 60
                    elif position == "bottom-left":
                        x, y = 40, 60
                    elif position == "bottom-center":
                        x, y = (pw - sig_w) / 2, 60
                    else:
                        x, y = pw - sig_w - 40, 60

                    c.drawImage(str(sig_path), x, y, width=sig_w, height=sig_h, preserveAspectRatio=True, mask="auto")

                    if date_text:
                        dt = datetime.now().strftime("%d/%m/%Y") if date_text == "auto" else date_text
                        c.setFont("Helvetica", 8)
                        c.setFillGray(0.4)
                        c.drawString(x, y - 12, dt)

                    c.save()
                    buf.seek(0)
                    overlay = PdfReader(buf)
                    pg.merge_page(overlay.pages[0])

                writer.add_page(pg)

            out = Path(output_path) if output_path else self._output_path(f"{pp.stem}_signed.pdf")
            out.parent.mkdir(parents=True, exist_ok=True)
            with open(str(out), "wb") as f:
                writer.write(f)
            logger.info(f"📄 Signature ajoutée: {out}")
            return {"success": True, "path": str(out), "filename": out.name}
        except Exception as e:
            logger.error(f"Erreur add_signature: {e}")
            return {"success": False, "error": str(e)}

    def list_pdf_fields(self, input_path: str) -> Dict[str, Any]:
        """Liste les champs remplissables d'un formulaire PDF."""
        try:
            from pypdf import PdfReader

            pp = Path(input_path)
            if not pp.exists():
                found = self._find_in_workspace(input_path)
                if found:
                    pp = found
                else:
                    return {"success": False, "error": f"Fichier non trouvé: {input_path}"}

            reader = PdfReader(str(pp))
            fields = {}
            if reader.get_fields():
                for name, field in reader.get_fields().items():
                    fields[name] = {
                        "type": str(field.get("/FT", "")),
                        "value": str(field.get("/V", "")),
                    }
            return {"success": True, "fields": fields, "count": len(fields)}
        except Exception as e:
            logger.error(f"Erreur list_pdf_fields: {e}")
            return {"success": False, "error": str(e)}

    def fill_pdf_form(
        self, input_path: str, fields: Union[Dict[str, str], str],
        output_filename: Optional[str] = None,
    ) -> Dict[str, Any]:
        """Remplit les champs d'un formulaire PDF."""
        try:
            from pypdf import PdfReader, PdfWriter

            fields = self._parse_json_arg(fields, {})
            pp = Path(input_path)
            if not pp.exists():
                found = self._find_in_workspace(input_path)
                if found:
                    pp = found
                else:
                    return {"success": False, "error": f"Fichier non trouvé: {input_path}"}

            reader = PdfReader(str(pp))
            writer = PdfWriter()
            writer.append(reader)

            for page_num in range(len(writer.pages)):
                writer.update_page_form_field_values(writer.pages[page_num], fields)

            out_name = output_filename or f"{pp.stem}_filled.pdf"
            out_path = self._output_path(out_name)
            with open(str(out_path), "wb") as f:
                writer.write(f)
            logger.info(f"📄 Formulaire PDF rempli: {out_path}")
            return {"success": True, "path": str(out_path), "filename": out_path.name, "fields_filled": len(fields)}
        except Exception as e:
            logger.error(f"Erreur fill_pdf_form: {e}")
            return {"success": False, "error": str(e)}

    def compare_documents(
        self, path_a: str, path_b: str
    ) -> Dict[str, Any]:
        """Compare deux documents et produit un rapport de différences."""
        try:
            import difflib

            # Lire les deux documents
            result_a = self.read_document(path_a)
            result_b = self.read_document(path_b)
            if not result_a.get("success"):
                return {"success": False, "error": f"Erreur lecture doc A: {result_a.get('error')}"}
            if not result_b.get("success"):
                return {"success": False, "error": f"Erreur lecture doc B: {result_b.get('error')}"}

            text_a = result_a.get("content", "").splitlines()
            text_b = result_b.get("content", "").splitlines()

            diff = list(difflib.unified_diff(text_a, text_b, fromfile=path_a, tofile=path_b, lineterm=""))
            additions = sum(1 for l in diff if l.startswith("+") and not l.startswith("+++"))
            deletions = sum(1 for l in diff if l.startswith("-") and not l.startswith("---"))

            diff_text = "\n".join(diff) if diff else "(Documents identiques)"

            # Sauvegarder le rapport
            out_name = "comparison_report.md"
            out_path = self._output_path(out_name)
            report = f"# Comparaison de documents\n\n"
            report += f"- **Document A** : {path_a}\n"
            report += f"- **Document B** : {path_b}\n"
            report += f"- **Ajouts** : {additions} lignes\n"
            report += f"- **Suppressions** : {deletions} lignes\n\n"
            report += f"## Diff\n\n```diff\n{diff_text}\n```\n"
            out_path.write_text(report, encoding="utf-8")

            return {
                "success": True, "path": str(out_path),
                "additions": additions, "deletions": deletions,
                "diff_preview": diff_text[:5000],
            }
        except Exception as e:
            logger.error(f"Erreur compare_documents: {e}")
            return {"success": False, "error": str(e)}

    def protect_pdf(
        self, input_path: str, password: str,
        output_path: Optional[str] = None,
        permissions: str = "read-only",
    ) -> Dict[str, Any]:
        """Protège un PDF par mot de passe."""
        try:
            from pypdf import PdfReader, PdfWriter
            from pypdf.constants import UserAccessPermissions

            pp = Path(input_path)
            if not pp.exists():
                found = self._find_in_workspace(input_path)
                if found:
                    pp = found
                else:
                    return {"success": False, "error": f"Fichier non trouvé: {input_path}"}

            reader = PdfReader(str(pp))
            writer = PdfWriter()
            writer.append(reader)

            # Map permissions
            perms = permissions.lower()
            if perms == "read-only":
                perm_flag = UserAccessPermissions.PRINT
            elif perms == "no-edit":
                perm_flag = UserAccessPermissions.PRINT | UserAccessPermissions.EXTRACT
            else:  # "full"
                perm_flag = UserAccessPermissions.all()

            writer.encrypt(
                user_password=password,
                owner_password=password + "_owner",
                permissions_flag=perm_flag,
            )

            out = Path(output_path) if output_path else self._output_path(f"{pp.stem}_protected.pdf")
            out.parent.mkdir(parents=True, exist_ok=True)
            with open(str(out), "wb") as f:
                writer.write(f)
            logger.info(f"🔒 PDF protégé: {out}")
            return {"success": True, "path": str(out), "filename": out.name, "permissions": perms}
        except Exception as e:
            logger.error(f"Erreur protect_pdf: {e}")
            return {"success": False, "error": str(e)}

    def image_to_document(
        self, image_path: str, output_format: str = "docx",
        language: str = "fra",
    ) -> Dict[str, Any]:
        """Convertit une image (ou PDF scanné) en document éditable via OCR."""
        try:
            pp = Path(image_path)
            if not pp.exists():
                found = self._find_in_workspace(image_path)
                if found:
                    pp = found
                else:
                    return {"success": False, "error": f"Fichier non trouvé: {image_path}"}

            text = ""
            if pp.suffix.lower() == ".pdf":
                # PDF scanné → OCR chaque page
                try:
                    from pdf2image import convert_from_path
                    import pytesseract
                    images = convert_from_path(str(pp))
                    for img in images:
                        text += pytesseract.image_to_string(img, lang=language) + "\n\n"
                except ImportError:
                    # Fallback: extraire le texte normalement
                    result = self.read_document(str(pp))
                    text = result.get("content", "")
            else:
                # Image simple
                try:
                    import pytesseract
                    from PIL import Image
                    img = Image.open(str(pp))
                    text = pytesseract.image_to_string(img, lang=language)
                except ImportError:
                    return {"success": False, "error": "pytesseract non installé. Installer tesseract-ocr."}

            if not text.strip():
                return {"success": False, "error": "Aucun texte détecté dans l'image."}

            stem = pp.stem
            fmt = output_format.lower().lstrip(".")
            if fmt == "pdf":
                return self.create_pdf(f"{stem}_ocr.pdf", title=f"OCR: {stem}", content=text)
            else:
                return self.create_docx(f"{stem}_ocr.docx", title=f"OCR: {stem}", content=text)

        except Exception as e:
            logger.error(f"Erreur image_to_document: {e}")
            return {"success": False, "error": str(e)}

    # ──────────────────────────────────────────────────────
    # HTML→PDF (via weasyprint) — P1.6
    # ──────────────────────────────────────────────────────

    def create_html_to_pdf(
        self, filename: str, html_content: str, css: str = ""
    ) -> Dict[str, Any]:
        """Convertit du HTML brut en PDF (WeasyPrint → Playwright fallback)."""
        filename = self._ensure_ext(filename, ".pdf")
        out_path = self._output_path(filename)
        full_html = html_content
        if css:
            full_html = f"<style>{css}</style>\n{html_content}"
        # Try WeasyPrint first
        try:
            import weasyprint
            weasyprint.HTML(string=full_html).write_pdf(str(out_path))
            logger.info(f"📄 HTML→PDF (weasyprint): {out_path}")
            return {"success": True, "path": str(out_path), "filename": filename}
        except Exception as wp_err:
            logger.warning(f"WeasyPrint indisponible ({wp_err}), fallback Playwright")
        # Fallback: Playwright Chromium
        try:
            from playwright.sync_api import sync_playwright
            with sync_playwright() as p:
                browser = p.chromium.launch(headless=True)
                page = browser.new_page()
                page.set_content(full_html, wait_until="networkidle")
                page.pdf(path=str(out_path), format="A4", print_background=True)
                browser.close()
            logger.info(f"📄 HTML→PDF (playwright): {out_path}")
            return {"success": True, "path": str(out_path), "filename": filename}
        except Exception as e:
            logger.error(f"Erreur html_to_pdf: {e}")
            return {"success": False, "error": str(e)}

    # ──────────────────────────────────────────────────────
    # Merge / Split PDF — P1.7
    # ──────────────────────────────────────────────────────

    def merge_pdfs(
        self, output_filename: str, input_paths: Union[List[str], str]
    ) -> Dict[str, Any]:
        """Fusionne plusieurs PDFs en un seul."""
        try:
            from pypdf import PdfWriter

            input_paths = self._parse_json_arg(input_paths, [])
            if not input_paths:
                return {"success": False, "error": "Aucun fichier à fusionner."}
            output_filename = self._ensure_ext(output_filename, ".pdf")
            out_path = self._output_path(output_filename)
            writer = PdfWriter()
            for p in input_paths:
                pp = Path(p)
                if not pp.exists():
                    return {"success": False, "error": f"Fichier non trouvé: {p}"}
                writer.append(str(pp))
            with open(str(out_path), "wb") as f:
                writer.write(f)
            logger.info(f"📄 PDF merge: {out_path} ({len(input_paths)} fichiers)")
            return {"success": True, "path": str(out_path), "filename": output_filename, "pages": len(writer.pages)}
        except Exception as e:
            logger.error(f"Erreur merge_pdfs: {e}")
            return {"success": False, "error": str(e)}

    def split_pdf(
        self, input_path: str, pages: str
    ) -> Dict[str, Any]:
        """Découpe un PDF. pages ex: '1-3,5,8-10'."""
        try:
            from pypdf import PdfReader, PdfWriter

            pp = Path(input_path)
            if not pp.exists():
                return {"success": False, "error": f"Fichier non trouvé: {input_path}"}
            reader = PdfReader(str(pp))
            total = len(reader.pages)
            # Parse page ranges
            indices = set()
            for part in pages.split(","):
                part = part.strip()
                if "-" in part:
                    a, b = part.split("-", 1)
                    indices.update(range(int(a) - 1, min(int(b), total)))
                else:
                    idx = int(part) - 1
                    if 0 <= idx < total:
                        indices.add(idx)
            if not indices:
                return {"success": False, "error": f"Aucune page valide dans '{pages}' (total: {total})"}
            writer = PdfWriter()
            for i in sorted(indices):
                writer.add_page(reader.pages[i])
            stem = pp.stem
            out_name = f"{stem}_pages_{pages.replace(',', '_').replace('-', 'to')}.pdf"
            out_path = self._output_path(out_name)
            with open(str(out_path), "wb") as f:
                writer.write(f)
            logger.info(f"📄 PDF split: {out_path} ({len(indices)} pages)")
            return {"success": True, "path": str(out_path), "filename": out_name, "pages_extracted": len(indices)}
        except Exception as e:
            logger.error(f"Erreur split_pdf: {e}")
            return {"success": False, "error": str(e)}

    # ──────────────────────────────────────────────────────
    # CSV — P1.8
    # ──────────────────────────────────────────────────────

    def create_csv(
        self,
        filename: str,
        headers: Union[List[str], str],
        rows: Union[List[List], str],
        delimiter: str = ",",
        encoding: str = "utf-8",
    ) -> Dict[str, Any]:
        """Crée un fichier CSV."""
        try:
            headers = self._parse_json_arg(headers, [])
            rows = self._parse_json_arg(rows, [])
            filename = self._ensure_ext(filename, ".csv")
            out_path = self._output_path(filename)
            with open(str(out_path), "w", newline="", encoding=encoding) as f:
                writer = csv.writer(f, delimiter=delimiter)
                if headers:
                    writer.writerow(headers)
                for row in rows:
                    writer.writerow(row)
            logger.info(f"📊 CSV créé: {out_path}")
            return {"success": True, "path": str(out_path), "filename": filename}
        except Exception as e:
            logger.error(f"Erreur create_csv: {e}")
            return {"success": False, "error": str(e)}

    # ──────────────────────────────────────────────────────
    # Conversion entre formats — P1.9
    # ──────────────────────────────────────────────────────

    def convert_document(
        self, input_path: str, output_format: str
    ) -> Dict[str, Any]:
        """
        Convertit un document d'un format à un autre.
        Conversions supportées: DOCX→PDF, DOCX→HTML, XLSX→CSV, HTML→PDF, MD→PDF, MD→DOCX, CSV→XLSX.
        """
        try:
            pp = Path(input_path)
            if not pp.exists():
                found = self._find_in_workspace(input_path)
                if found:
                    pp = found
                else:
                    return {"success": False, "error": f"Fichier non trouvé: {input_path}"}

            src_ext = pp.suffix.lower()
            dst = output_format.lower().lstrip(".")
            stem = pp.stem

            if src_ext == ".docx" and dst == "pdf":
                import mammoth
                with open(str(pp), "rb") as f:
                    result = mammoth.convert_to_html(f)
                html = result.value
                return self.create_html_to_pdf(f"{stem}.pdf", html)

            elif src_ext == ".docx" and dst == "html":
                import mammoth
                with open(str(pp), "rb") as f:
                    result = mammoth.convert_to_html(f)
                out_name = f"{stem}.html"
                out_path = self._output_path(out_name)
                out_path.write_text(result.value, encoding="utf-8")
                return {"success": True, "path": str(out_path), "filename": out_name}

            elif src_ext == ".xlsx" and dst == "csv":
                import openpyxl
                wb = openpyxl.load_workbook(str(pp), data_only=True)
                ws = wb.active
                out_name = f"{stem}.csv"
                out_path = self._output_path(out_name)
                with open(str(out_path), "w", newline="", encoding="utf-8") as f:
                    writer = csv.writer(f)
                    for row in ws.iter_rows(values_only=True):
                        writer.writerow([c if c is not None else "" for c in row])
                return {"success": True, "path": str(out_path), "filename": out_name}

            elif src_ext == ".html" and dst == "pdf":
                html = pp.read_text(encoding="utf-8")
                return self.create_html_to_pdf(f"{stem}.pdf", html)

            elif src_ext in (".md", ".markdown") and dst == "pdf":
                md_content = pp.read_text(encoding="utf-8")
                return self.create_pdf(f"{stem}.pdf", title=stem, content=md_content)

            elif src_ext in (".md", ".markdown") and dst == "docx":
                md_content = pp.read_text(encoding="utf-8")
                return self.create_docx(f"{stem}.docx", title=stem, content=md_content)

            elif src_ext == ".csv" and dst == "xlsx":
                with open(str(pp), "r", encoding="utf-8") as f:
                    reader = csv.reader(f)
                    all_rows = list(reader)
                if all_rows:
                    headers = all_rows[0]
                    data_rows = all_rows[1:]
                else:
                    headers, data_rows = [], []
                return self.create_xlsx(f"{stem}.xlsx", [{"name": stem, "headers": headers, "rows": data_rows}])

            else:
                return {"success": False, "error": f"Conversion {src_ext}→{dst} non supportée."}

        except Exception as e:
            logger.error(f"Erreur convert_document: {e}")
            return {"success": False, "error": str(e)}

    # ──────────────────────────────────────────────────────
    # Édition DOCX — P1B.1
    # ──────────────────────────────────────────────────────

    def edit_docx(
        self, input_path: str, operations: Union[List[Dict], str], output_path: Optional[str] = None
    ) -> Dict[str, Any]:
        """Édite un fichier DOCX existant via une liste d'opérations."""
        try:
            from docx import Document
            from docx.shared import Cm

            operations = self._parse_json_arg(operations, [])
            pp = Path(input_path)
            if not pp.exists():
                found = self._find_in_workspace(input_path)
                if found:
                    pp = found
                else:
                    return {"success": False, "error": f"Fichier non trouvé: {input_path}"}

            doc = Document(str(pp))
            ops_done = 0

            for op_def in operations:
                op = op_def.get("op", "")
                if op == "replace_text":
                    find = op_def.get("find", "")
                    replace = op_def.get("replace", "")
                    for para in doc.paragraphs:
                        if find in para.text:
                            for run in para.runs:
                                if find in run.text:
                                    run.text = run.text.replace(find, replace)
                                    ops_done += 1
                    # Also search in tables
                    for table in doc.tables:
                        for row in table.rows:
                            for cell in row.cells:
                                for para in cell.paragraphs:
                                    if find in para.text:
                                        for run in para.runs:
                                            if find in run.text:
                                                run.text = run.text.replace(find, replace)
                                                ops_done += 1

                elif op == "add_paragraph":
                    text = op_def.get("text", "")
                    style = op_def.get("style", "Normal")
                    p = doc.add_paragraph(text, style=style)
                    ops_done += 1

                elif op == "delete_paragraph":
                    containing = op_def.get("containing", "")
                    for para in doc.paragraphs:
                        if containing and containing in para.text:
                            p_element = para._element
                            p_element.getparent().remove(p_element)
                            ops_done += 1
                            break

                elif op == "add_image":
                    img_path = op_def.get("path", "")
                    w_cm = op_def.get("width_cm", 12)
                    if img_path and Path(img_path).exists():
                        doc.add_picture(img_path, width=Cm(w_cm))
                        ops_done += 1

                elif op == "set_header":
                    text = op_def.get("text", "")
                    section = doc.sections[0]
                    header = section.header
                    header.is_linked_to_previous = False
                    hp = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
                    hp.text = text
                    ops_done += 1

                elif op == "set_footer":
                    text = op_def.get("text", "")
                    section = doc.sections[0]
                    footer = section.footer
                    footer.is_linked_to_previous = False
                    fp = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
                    fp.text = text
                    ops_done += 1

                elif op == "replace_in_table":
                    ti = op_def.get("table_index", 0)
                    ri = op_def.get("row", 0)
                    ci = op_def.get("col", 0)
                    val = str(op_def.get("value", ""))
                    if ti < len(doc.tables):
                        table = doc.tables[ti]
                        if ri < len(table.rows) and ci < len(table.rows[ri].cells):
                            table.rows[ri].cells[ci].text = val
                            ops_done += 1

            out = Path(output_path) if output_path else pp
            out.parent.mkdir(parents=True, exist_ok=True)
            doc.save(str(out))
            logger.info(f"📝 DOCX édité: {out} ({ops_done} opérations)")
            return {"success": True, "path": str(out), "operations_applied": ops_done}
        except Exception as e:
            logger.error(f"Erreur edit_docx: {e}")
            return {"success": False, "error": str(e)}

    # ──────────────────────────────────────────────────────
    # Édition XLSX — P1B.2
    # ──────────────────────────────────────────────────────

    def edit_xlsx(
        self, input_path: str, operations: Union[List[Dict], str], output_path: Optional[str] = None
    ) -> Dict[str, Any]:
        """Édite un fichier XLSX existant via une liste d'opérations."""
        try:
            import openpyxl

            operations = self._parse_json_arg(operations, [])
            pp = Path(input_path)
            if not pp.exists():
                found = self._find_in_workspace(input_path)
                if found:
                    pp = found
                else:
                    return {"success": False, "error": f"Fichier non trouvé: {input_path}"}

            wb = openpyxl.load_workbook(str(pp))
            ops_done = 0

            for op_def in operations:
                op = op_def.get("op", "")
                sheet_name = op_def.get("sheet", wb.sheetnames[0] if wb.sheetnames else "Sheet1")

                if op == "set_cell":
                    ws = wb[sheet_name] if sheet_name in wb.sheetnames else wb.active
                    cell_ref = op_def.get("cell", "A1")
                    ws[cell_ref] = op_def.get("value", "")
                    ops_done += 1

                elif op == "set_formula":
                    ws = wb[sheet_name] if sheet_name in wb.sheetnames else wb.active
                    cell_ref = op_def.get("cell", "A1")
                    ws[cell_ref] = op_def.get("formula", "")
                    ops_done += 1

                elif op == "add_row":
                    ws = wb[sheet_name] if sheet_name in wb.sheetnames else wb.active
                    values = op_def.get("values", [])
                    ws.append(values)
                    ops_done += 1

                elif op == "delete_row":
                    ws = wb[sheet_name] if sheet_name in wb.sheetnames else wb.active
                    row_num = op_def.get("row", 1)
                    ws.delete_rows(row_num)
                    ops_done += 1

                elif op == "add_sheet":
                    name = op_def.get("name", "Nouvelle feuille")
                    wb.create_sheet(title=name)
                    ops_done += 1

                elif op == "rename_sheet":
                    old = op_def.get("old_name", "")
                    new = op_def.get("new_name", "")
                    if old in wb.sheetnames:
                        wb[old].title = new
                        ops_done += 1

            out = Path(output_path) if output_path else pp
            out.parent.mkdir(parents=True, exist_ok=True)
            wb.save(str(out))
            logger.info(f"📊 XLSX édité: {out} ({ops_done} opérations)")
            return {"success": True, "path": str(out), "operations_applied": ops_done}
        except Exception as e:
            logger.error(f"Erreur edit_xlsx: {e}")
            return {"success": False, "error": str(e)}

    # ──────────────────────────────────────────────────────
    # Édition PPTX — P1B.3
    # ──────────────────────────────────────────────────────

    def edit_pptx(
        self, input_path: str, operations: Union[List[Dict], str], output_path: Optional[str] = None
    ) -> Dict[str, Any]:
        """Édite un fichier PPTX existant via une liste d'opérations."""
        try:
            from pptx import Presentation
            from pptx.util import Cm

            operations = self._parse_json_arg(operations, [])
            pp = Path(input_path)
            if not pp.exists():
                found = self._find_in_workspace(input_path)
                if found:
                    pp = found
                else:
                    return {"success": False, "error": f"Fichier non trouvé: {input_path}"}

            prs = Presentation(str(pp))
            ops_done = 0

            for op_def in operations:
                op = op_def.get("op", "")
                if op == "replace_text":
                    slide_idx = op_def.get("slide", 1) - 1
                    find = op_def.get("find", "")
                    replace = op_def.get("replace", "")
                    if 0 <= slide_idx < len(prs.slides):
                        slide = prs.slides[slide_idx]
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    for run in para.runs:
                                        if find in run.text:
                                            run.text = run.text.replace(find, replace)
                                            ops_done += 1

                elif op == "add_slide":
                    t = op_def.get("title", "")
                    c = op_def.get("content", "")
                    blank = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank)
                    from pptx.util import Inches, Pt
                    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
                    p = tb.text_frame.paragraphs[0]
                    r = p.add_run()
                    r.text = t
                    r.font.size = Pt(24)
                    r.font.bold = True
                    if c:
                        tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(5))
                        p2 = tb2.text_frame.paragraphs[0]
                        p2.add_run().text = c
                    ops_done += 1

                elif op == "delete_slide":
                    slide_idx = op_def.get("slide", 1) - 1
                    if 0 <= slide_idx < len(prs.slides):
                        rId = prs.slides._sldIdLst[slide_idx].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                        prs.part.drop_rel(rId)
                        del prs.slides._sldIdLst[slide_idx]
                        ops_done += 1

                elif op == "add_image":
                    slide_idx = op_def.get("slide", 1) - 1
                    img_path = op_def.get("path", "")
                    if 0 <= slide_idx < len(prs.slides) and img_path and Path(img_path).exists():
                        slide = prs.slides[slide_idx]
                        left = Cm(op_def.get("left_cm", 2))
                        top = Cm(op_def.get("top_cm", 3))
                        width = Cm(op_def.get("width_cm", 10))
                        slide.shapes.add_picture(img_path, left, top, width)
                        ops_done += 1

            out = Path(output_path) if output_path else pp
            out.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(out))
            logger.info(f"🎞️ PPTX édité: {out} ({ops_done} opérations)")
            return {"success": True, "path": str(out), "operations_applied": ops_done}
        except Exception as e:
            logger.error(f"Erreur edit_pptx: {e}")
            return {"success": False, "error": str(e)}

    # ──────────────────────────────────────────────────────
    # Annoter PDF — P1B.4
    # ──────────────────────────────────────────────────────

    def annotate_pdf(
        self, input_path: str, annotations: Union[List[Dict], str], output_path: Optional[str] = None
    ) -> Dict[str, Any]:
        """Annote un PDF existant (texte, surlignage, tampon)."""
        try:
            from pypdf import PdfReader, PdfWriter
            from reportlab.lib.colors import HexColor
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.units import cm
            from reportlab.pdfgen.canvas import Canvas

            annotations = self._parse_json_arg(annotations, [])
            pp = Path(input_path)
            if not pp.exists():
                found = self._find_in_workspace(input_path)
                if found:
                    pp = found
                else:
                    return {"success": False, "error": f"Fichier non trouvé: {input_path}"}

            reader = PdfReader(str(pp))
            writer = PdfWriter()
            total_pages = len(reader.pages)

            # Group annotations by page
            page_annots: Dict[int, list] = {}
            for ann in annotations:
                page = ann.get("page", 0)
                if page == -1:
                    page = total_pages - 1
                page_annots.setdefault(page, []).append(ann)

            for page_idx in range(total_pages):
                page = reader.pages[page_idx]
                if page_idx not in page_annots:
                    writer.add_page(page)
                    continue

                # Create overlay
                mbox = page.mediabox
                pw, ph = float(mbox.width), float(mbox.height)
                overlay_buf = io.BytesIO()
                c = Canvas(overlay_buf, pagesize=(pw, ph))

                for ann in page_annots[page_idx]:
                    atype = ann.get("type", "text")
                    if atype == "text":
                        x = float(ann.get("x", 100))
                        y = float(ann.get("y", 500))
                        text = ann.get("text", "")
                        fs = int(ann.get("font_size", 12))
                        color = ann.get("color", "black")
                        try:
                            c.setFillColor(HexColor(color))
                        except Exception:
                            c.setFillColor(HexColor("#000000"))
                        c.setFont("Helvetica", fs)
                        c.drawString(x, y, text)

                    elif atype == "highlight":
                        x1 = float(ann.get("x1", 50))
                        y1 = float(ann.get("y1", 600))
                        x2 = float(ann.get("x2", 400))
                        y2 = float(ann.get("y2", 620))
                        c.setFillColor(HexColor("#FFFF00"))
                        c.setFillAlpha(0.3)
                        c.rect(x1, y1, x2 - x1, y2 - y1, fill=1, stroke=0)
                        c.setFillAlpha(1.0)

                    elif atype == "stamp":
                        text = ann.get("text", "VALIDÉ")
                        pos = ann.get("position", "top-right")
                        c.setFillColor(HexColor("#CC0000"))
                        c.setFillAlpha(0.5)
                        c.setFont("Helvetica-Bold", 36)
                        if pos == "top-right":
                            c.drawRightString(pw - 40, ph - 60, text)
                        elif pos == "center":
                            c.drawCentredString(pw / 2, ph / 2, text)
                        else:
                            c.drawString(40, ph - 60, text)
                        c.setFillAlpha(1.0)

                c.save()
                overlay_buf.seek(0)
                overlay_reader = PdfReader(overlay_buf)
                page.merge_page(overlay_reader.pages[0])
                writer.add_page(page)

            out = Path(output_path) if output_path else self._output_path(f"{pp.stem}_annotated.pdf")
            out.parent.mkdir(parents=True, exist_ok=True)
            with open(str(out), "wb") as f:
                writer.write(f)
            logger.info(f"📄 PDF annoté: {out}")
            return {"success": True, "path": str(out), "filename": out.name, "annotations_count": len(annotations)}
        except Exception as e:
            logger.error(f"Erreur annotate_pdf: {e}")
            return {"success": False, "error": str(e)}

    # ──────────────────────────────────────────────────────
    # P5 — Formats supplémentaires
    # ──────────────────────────────────────────────────────

    def create_markdown(
        self,
        filename: str,
        title: str,
        content: str,
    ) -> Dict[str, Any]:
        """Crée un fichier Markdown (.md)."""
        try:
            filename = self._ensure_ext(filename, ".md")
            out_path = self._output_path(filename)
            md = f"# {title}\n\n{content}\n"
            with open(str(out_path), "w", encoding="utf-8") as f:
                f.write(md)
            logger.info(f"📝 Markdown créé: {out_path}")
            return {"success": True, "path": str(out_path), "filename": filename}
        except Exception as e:
            logger.error(f"Erreur create_markdown: {e}")
            return {"success": False, "error": str(e)}

    def create_html(
        self,
        filename: str,
        title: str,
        content: str,
        css: str = "",
        template: str = "default",
    ) -> Dict[str, Any]:
        """
        Crée un fichier HTML standalone avec CSS embarqué.
        Templates: default, report, email, print.
        """
        try:
            filename = self._ensure_ext(filename, ".html")
            out_path = self._output_path(filename)

            _CSS = {
                "default": "body{font-family:system-ui,sans-serif;max-width:800px;margin:2rem auto;padding:0 1rem;color:#1e293b;line-height:1.6}h1{color:#1a1a2e;border-bottom:2px solid #6366f1;padding-bottom:.5rem}h2{color:#2d3561}a{color:#6366f1}code{background:#f1f5f9;padding:2px 6px;border-radius:3px;font-size:.9em}pre{background:#1e293b;color:#e2e8f0;padding:1rem;border-radius:6px;overflow-x:auto}table{border-collapse:collapse;width:100%}th,td{border:1px solid #cbd5e1;padding:.5rem;text-align:left}th{background:#f1f5f9;font-weight:600}",
                "report": "body{font-family:Georgia,serif;max-width:900px;margin:2rem auto;padding:0 2rem;color:#222;line-height:1.8}h1{text-align:center;color:#1a1a2e;font-size:2em}h2{color:#2d3561;border-bottom:1px solid #ccc;padding-bottom:.3rem}table{border-collapse:collapse;width:100%}th,td{border:1px solid #999;padding:.4rem}th{background:#eee}",
                "email": "body{font-family:Arial,sans-serif;max-width:600px;margin:0 auto;padding:1rem;color:#333;line-height:1.5}h1{font-size:1.4em;color:#1a1a2e}a{color:#6366f1}",
                "print": "@media print{body{font-size:12pt}}body{font-family:Georgia,serif;max-width:700px;margin:2rem auto;color:#000;line-height:1.6}h1{font-size:1.8em}h2{font-size:1.4em}table{border-collapse:collapse;width:100%}th,td{border:1px solid #000;padding:.3rem}",
            }
            base_css = _CSS.get(template, _CSS["default"])
            if css:
                base_css += "\n" + css

            # Convert markdown content to basic HTML
            html_body = self._md_to_html(content)

            html = f"""<!DOCTYPE html>
<html lang="fr">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title}</title>
<style>{base_css}</style>
</head>
<body>
<h1>{title}</h1>
{html_body}
</body>
</html>"""
            with open(str(out_path), "w", encoding="utf-8") as f:
                f.write(html)
            logger.info(f"🌐 HTML créé: {out_path}")
            return {"success": True, "path": str(out_path), "filename": filename}
        except Exception as e:
            logger.error(f"Erreur create_html: {e}")
            return {"success": False, "error": str(e)}

    def create_email_html(
        self,
        filename: str,
        subject: str,
        body: str,
        sender_name: str = "Lumena",
        sender_logo_url: str = "",
    ) -> Dict[str, Any]:
        """
        Crée un HTML compatible email (CSS inline, tables pour layout).
        Compatible Gmail/Outlook.
        """
        try:
            filename = self._ensure_ext(filename, ".html")
            out_path = self._output_path(filename)

            body_html = self._md_to_html(body)
            logo_block = ""
            if sender_logo_url:
                logo_block = f'<img src="{sender_logo_url}" alt="{sender_name}" style="max-width:150px;height:auto;margin-bottom:16px;">'

            html = f"""<!DOCTYPE html>
<html lang="fr">
<head><meta charset="UTF-8"><title>{subject}</title></head>
<body style="margin:0;padding:0;background-color:#f4f4f7;font-family:Arial,sans-serif;">
<table role="presentation" width="100%" cellpadding="0" cellspacing="0" style="background-color:#f4f4f7;">
<tr><td align="center" style="padding:40px 0;">
<table role="presentation" width="600" cellpadding="0" cellspacing="0" style="background-color:#ffffff;border-radius:8px;overflow:hidden;box-shadow:0 2px 8px rgba(0,0,0,0.08);">
<tr><td style="background-color:#1a1a2e;padding:24px 32px;text-align:center;">
{logo_block}
<h1 style="color:#ffffff;margin:0;font-size:22px;font-weight:600;">{subject}</h1>
</td></tr>
<tr><td style="padding:32px;color:#333333;font-size:15px;line-height:1.6;">
{body_html}
</td></tr>
<tr><td style="background-color:#f8f9fa;padding:16px 32px;text-align:center;color:#888;font-size:12px;">
Envoyé par {sender_name}
</td></tr>
</table>
</td></tr>
</table>
</body>
</html>"""
            with open(str(out_path), "w", encoding="utf-8") as f:
                f.write(html)
            logger.info(f"📧 Email HTML créé: {out_path}")
            return {"success": True, "path": str(out_path), "filename": filename}
        except Exception as e:
            logger.error(f"Erreur create_email_html: {e}")
            return {"success": False, "error": str(e)}

    def create_ics(
        self,
        filename: str,
        events: Any,
    ) -> Dict[str, Any]:
        """
        Crée un fichier iCalendar (.ics) RFC 5545.
        events: [{title, start, end, location, description, attendees[]}]
        Dates au format ISO 8601 ou "YYYY-MM-DD HH:MM".
        """
        try:
            events_data = self._parse_json_arg(events, [])
            if isinstance(events_data, dict):
                events_data = [events_data]
            if not events_data:
                return {"success": False, "error": "Aucun événement fourni."}

            filename = self._ensure_ext(filename, ".ics")
            out_path = self._output_path(filename)

            def _dt(s: str) -> str:
                """Convertit date string en format ICS YYYYMMDDTHHMMSS."""
                s = s.strip().replace("-", "").replace(":", "").replace(" ", "T")
                if "T" not in s:
                    s += "T000000"
                if len(s) < 15:
                    s = s.ljust(15, "0")
                return s[:15]

            lines = ["BEGIN:VCALENDAR", "VERSION:2.0", "PRODID:-//Lumena//Documents//FR"]
            for ev in events_data:
                lines.append("BEGIN:VEVENT")
                lines.append(f"SUMMARY:{ev.get('title', 'Sans titre')}")
                lines.append(f"DTSTART:{_dt(ev.get('start', ''))}")
                lines.append(f"DTEND:{_dt(ev.get('end', ev.get('start', '')))}")
                if ev.get("location"):
                    lines.append(f"LOCATION:{ev['location']}")
                if ev.get("description"):
                    lines.append(f"DESCRIPTION:{ev['description']}")
                for att in ev.get("attendees", []):
                    lines.append(f"ATTENDEE:mailto:{att}")
                lines.append("END:VEVENT")
            lines.append("END:VCALENDAR")

            with open(str(out_path), "w", encoding="utf-8") as f:
                f.write("\r\n".join(lines) + "\r\n")
            logger.info(f"📅 ICS créé: {out_path}")
            return {"success": True, "path": str(out_path), "filename": filename}
        except Exception as e:
            logger.error(f"Erreur create_ics: {e}")
            return {"success": False, "error": str(e)}

    def create_vcard(
        self,
        filename: str,
        contacts: Any,
    ) -> Dict[str, Any]:
        """
        Crée un fichier vCard 3.0 (.vcf).
        contacts: [{name, company, phone, email, address, website, photo_url}]
        """
        try:
            contacts_data = self._parse_json_arg(contacts, [])
            if isinstance(contacts_data, dict):
                contacts_data = [contacts_data]
            if not contacts_data:
                return {"success": False, "error": "Aucun contact fourni."}

            filename = self._ensure_ext(filename, ".vcf")
            out_path = self._output_path(filename)

            lines = []
            for c in contacts_data:
                name = c.get("name", "Sans nom")
                parts = name.split(" ", 1)
                first = parts[0]
                last = parts[1] if len(parts) > 1 else ""
                lines.append("BEGIN:VCARD")
                lines.append("VERSION:3.0")
                lines.append(f"FN:{name}")
                lines.append(f"N:{last};{first};;;")
                if c.get("company"):
                    lines.append(f"ORG:{c['company']}")
                if c.get("phone"):
                    lines.append(f"TEL;TYPE=CELL:{c['phone']}")
                if c.get("email"):
                    lines.append(f"EMAIL:{c['email']}")
                if c.get("address"):
                    lines.append(f"ADR;TYPE=HOME:;;{c['address']};;;;")
                if c.get("website"):
                    lines.append(f"URL:{c['website']}")
                if c.get("photo_url"):
                    lines.append(f"PHOTO;VALUE=uri:{c['photo_url']}")
                lines.append("END:VCARD")

            with open(str(out_path), "w", encoding="utf-8") as f:
                f.write("\r\n".join(lines) + "\r\n")
            logger.info(f"👤 vCard créé: {out_path}")
            return {"success": True, "path": str(out_path), "filename": filename}
        except Exception as e:
            logger.error(f"Erreur create_vcard: {e}")
            return {"success": False, "error": str(e)}

    @staticmethod
    def _md_to_html(content: str) -> str:
        """Convertit un contenu markdown basique en HTML."""
        import html as _html
        lines = content.split("\n")
        out = []
        in_list = False
        for line in lines:
            s = line.strip()
            if not s:
                if in_list:
                    out.append("</ul>")
                    in_list = False
                out.append("<br>")
                continue
            if s.startswith("### "):
                if in_list:
                    out.append("</ul>")
                    in_list = False
                out.append(f"<h3>{_html.escape(s[4:])}</h3>")
            elif s.startswith("## "):
                if in_list:
                    out.append("</ul>")
                    in_list = False
                out.append(f"<h2>{_html.escape(s[3:])}</h2>")
            elif s.startswith("# "):
                if in_list:
                    out.append("</ul>")
                    in_list = False
                out.append(f"<h2>{_html.escape(s[2:])}</h2>")
            elif s.startswith(("- ", "* ")):
                if not in_list:
                    out.append("<ul>")
                    in_list = True
                text = s[2:]
                text = re.sub(r"\*\*(.*?)\*\*", r"<strong>\1</strong>", text)
                text = re.sub(r"\*(.*?)\*", r"<em>\1</em>", text)
                out.append(f"<li>{text}</li>")
            elif re.match(r"^-{3,}$", s):
                if in_list:
                    out.append("</ul>")
                    in_list = False
                out.append("<hr>")
            else:
                if in_list:
                    out.append("</ul>")
                    in_list = False
                text = _html.escape(s)
                text = re.sub(r"\*\*(.*?)\*\*", r"<strong>\1</strong>", text)
                text = re.sub(r"\*(.*?)\*", r"<em>\1</em>", text)
                text = re.sub(r"`(.*?)`", r"<code>\1</code>", text)
                out.append(f"<p>{text}</p>")
        if in_list:
            out.append("</ul>")
        return "\n".join(out)

    # ──────────────────────────────────────────────────────
    # P6 — Batch et automatisation
    # ──────────────────────────────────────────────────────

    def create_batch_documents(
        self,
        template_name: str,
        data_source: Any,
        output_format: str = "pdf",
    ) -> Dict[str, Any]:
        """
        Génère un lot de documents à partir d'un template et d'une source de données.
        data_source: chemin CSV/XLSX/JSON, ou liste de dicts.
        """
        try:
            from jinja2 import Environment, FileSystemLoader

            data = self._parse_json_arg(data_source, None)

            # Si c'est un chemin fichier, lire les données
            if isinstance(data, str) or data is None:
                src_path = Path(str(data_source)) if data is None else Path(str(data))
                if src_path.exists():
                    ext = src_path.suffix.lower()
                    if ext == ".csv":
                        with open(str(src_path), "r", encoding="utf-8") as f:
                            reader = csv.DictReader(f)
                            data = list(reader)
                    elif ext in (".xlsx", ".xls"):
                        import openpyxl
                        wb = openpyxl.load_workbook(str(src_path), data_only=True)
                        ws = wb.active
                        headers = [str(c.value or "") for c in ws[1]]
                        data = []
                        for row in ws.iter_rows(min_row=2, values_only=True):
                            data.append({h: (str(v) if v is not None else "") for h, v in zip(headers, row)})
                    elif ext == ".json":
                        with open(str(src_path), "r", encoding="utf-8") as f:
                            data = json.load(f)
                    else:
                        return {"success": False, "error": f"Format source non supporté: {ext}"}
                else:
                    return {"success": False, "error": f"Fichier source introuvable: {data_source}"}

            if not isinstance(data, list) or not data:
                return {"success": False, "error": "data_source doit être une liste non vide de dicts."}

            # Chercher le template
            from src.utils.paths import TEMPLATES_DIR
            templates_dir = TEMPLATES_DIR
            custom_dir = templates_dir / "custom"
            env = Environment(loader=FileSystemLoader([str(custom_dir), str(templates_dir)]))

            tpl_file = f"{template_name}.html.j2"
            try:
                template = env.get_template(tpl_file)
            except Exception:
                return {"success": False, "error": f"Template '{template_name}' introuvable."}

            results = []
            for i, row_data in enumerate(data):
                html_content = template.render(**row_data, index=i + 1)

                # Générer un nom de fichier unique
                row_name = row_data.get("name", row_data.get("nom", row_data.get("title", f"doc_{i+1}")))
                safe_name = re.sub(r"[^\w\-]", "_", str(row_name))
                fname = f"{template_name}_{safe_name}"

                if output_format == "pdf":
                    result = self.create_html_to_pdf(filename=f"{fname}.pdf", html_content=html_content, title=str(row_name))
                elif output_format == "html":
                    out_path = self._output_path(f"{fname}.html")
                    with open(str(out_path), "w", encoding="utf-8") as f:
                        f.write(html_content)
                    result = {"success": True, "path": str(out_path), "filename": f"{fname}.html"}
                else:
                    result = {"success": False, "error": f"Format non supporté: {output_format}"}

                results.append(result)

            succeeded = sum(1 for r in results if r.get("success"))
            paths = [r["path"] for r in results if r.get("success")]
            return {
                "success": succeeded > 0,
                "total": len(data),
                "succeeded": succeeded,
                "failed": len(data) - succeeded,
                "paths": paths,
            }
        except Exception as e:
            logger.error(f"Erreur create_batch_documents: {e}")
            return {"success": False, "error": str(e)}

    def zip_documents(
        self,
        output_filename: str,
        paths: Any,
    ) -> Dict[str, Any]:
        """
        Crée un fichier ZIP contenant les documents listés.
        paths: liste de chemins fichiers.
        """
        try:
            paths_list = self._parse_json_arg(paths, [])
            if not paths_list:
                return {"success": False, "error": "Aucun chemin fourni."}

            output_filename = self._ensure_ext(output_filename, ".zip")
            out_path = self._output_path(output_filename)

            with zipfile.ZipFile(str(out_path), "w", zipfile.ZIP_DEFLATED) as zf:
                for p in paths_list:
                    fp = Path(p)
                    if fp.exists() and fp.is_file():
                        zf.write(str(fp), fp.name)

            entries = len([p for p in paths_list if Path(p).exists()])
            logger.info(f"📦 ZIP créé: {out_path} ({entries} fichiers)")
            return {
                "success": True,
                "path": str(out_path),
                "filename": output_filename,
                "entries": entries,
            }
        except Exception as e:
            logger.error(f"Erreur zip_documents: {e}")
            return {"success": False, "error": str(e)}

    def assemble_document(
        self,
        output_filename: str,
        parts: Any,
    ) -> Dict[str, Any]:
        """
        Assemble un document composite à partir de blocs.
        parts: liste de dicts [{type: "pdf"|"section"|"text", content/path: ...}]
        - type "pdf": merge un PDF existant (path requis)
        - type "section": génère une section PDF (title + content)
        - type "text": ajoute du texte brut
        """
        try:
            from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
            from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
            from reportlab.lib.units import cm
            from reportlab.lib import colors
            from reportlab.platypus import (
                HRFlowable, Paragraph, SimpleDocTemplate, Spacer, PageBreak,
            )
            from pypdf import PdfReader, PdfWriter

            parts_data = self._parse_json_arg(parts, [])
            if not parts_data:
                return {"success": False, "error": "Aucune partie fournie."}

            output_filename = self._ensure_ext(output_filename, ".pdf")
            out_path = self._output_path(output_filename)

            # Step 1: generate intermediate PDFs for non-PDF parts
            import tempfile
            temp_pdfs: List[str] = []

            for idx, part in enumerate(parts_data):
                part_type = part.get("type", "text")

                if part_type == "pdf":
                    pdf_path = part.get("path", "")
                    if pdf_path and Path(pdf_path).exists():
                        temp_pdfs.append(pdf_path)
                    continue

                # Generate a temp PDF for section/text parts
                content = part.get("content", "")
                title = part.get("title", "")

                styles = getSampleStyleSheet()
                story = []

                if title:
                    style_title = ParagraphStyle(
                        f"ATitle{idx}", parent=styles["Title"],
                        fontSize=18, textColor=colors.HexColor("#1a1a2e"),
                        spaceAfter=12, alignment=TA_CENTER,
                    )
                    story.append(Paragraph(self._md_inline_rl(title), style_title))
                    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#6366f1"), spaceAfter=10))

                if content:
                    style_body = ParagraphStyle(
                        f"ABody{idx}", parent=styles["Normal"],
                        fontSize=11, leading=16, spaceAfter=5, alignment=TA_JUSTIFY,
                    )
                    for line in content.split("\n"):
                        s = line.strip()
                        if not s:
                            story.append(Spacer(1, 4))
                        else:
                            story.append(Paragraph(self._md_inline_rl(s), style_body))

                if story:
                    fd, tmp_path = tempfile.mkstemp(suffix=".pdf")
                    os.close(fd)
                    tmp_doc = SimpleDocTemplate(tmp_path, pagesize=(595.28, 841.89))  # A4
                    tmp_doc.build(story)
                    temp_pdfs.append(tmp_path)

            # Step 2: merge all PDFs
            writer = PdfWriter()
            generated_temps = []
            for p in temp_pdfs:
                reader = PdfReader(p)
                for page in reader.pages:
                    writer.add_page(page)
                # Track temp files for cleanup (not user-provided ones)
                if not any(part.get("path") == p for part in parts_data):
                    generated_temps.append(p)

            with open(str(out_path), "wb") as f:
                writer.write(f)

            # Cleanup temp files
            for tp in generated_temps:
                try:
                    os.unlink(tp)
                except OSError:
                    pass

            logger.info(f"📄 Document composite créé: {out_path} ({len(temp_pdfs)} parties)")
            return {
                "success": True,
                "path": str(out_path),
                "filename": output_filename,
                "parts_count": len(temp_pdfs),
            }
        except Exception as e:
            logger.error(f"Erreur assemble_document: {e}")
            return {"success": False, "error": str(e)}

    # ──────────────────────────────────────────────────────
    # Lire un document existant (DOCX / XLSX / PPTX / PDF)
    # ──────────────────────────────────────────────────────

    def read_document(self, path: str) -> Dict[str, Any]:
        """
        Lit le contenu textuel d'un document Office ou PDF.
        Retourne {"success": True, "content": str, "type": str, "path": str}
        """
        try:
            p = Path(path)
            if not p.is_absolute():
                # Chercher dans le workspace
                p = self._find_in_workspace(path)
                if p is None:
                    return {"success": False, "error": f"Fichier non trouvé: {path}"}

            if not p.exists():
                return {"success": False, "error": f"Fichier non trouvé: {p}"}

            ext = p.suffix.lower()

            if ext == ".docx":
                return self._read_docx(p)
            elif ext == ".xlsx":
                return self._read_xlsx(p)
            elif ext == ".pptx":
                return self._read_pptx(p)
            elif ext == ".pdf":
                return self._read_pdf(p)
            else:
                # Lecture texte brute
                content = p.read_text(encoding="utf-8", errors="replace")
                return {"success": True, "content": content, "type": ext, "path": str(p)}

        except Exception as e:
            logger.error(f"Erreur read_document: {e}")
            return {"success": False, "error": str(e)}

    def _find_in_workspace(self, filename: str) -> Optional[Path]:
        """Recherche récursive d'un fichier par nom dans le workspace."""
        name = Path(filename).name
        for candidate in self.workspace_root.rglob(name):
            if candidate.is_file():
                return candidate
        return None

    def _read_docx(self, p: Path) -> Dict[str, Any]:
        from docx import Document
        doc = Document(str(p))
        lines = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                lines.append("")
                continue
            style_name = (para.style.name or "").lower()
            if "heading 1" in style_name or style_name == "title":
                lines.append(f"# {text}")
            elif "heading 2" in style_name:
                lines.append(f"## {text}")
            elif "heading 3" in style_name:
                lines.append(f"### {text}")
            elif "list bullet" in style_name:
                lines.append(f"- {text}")
            elif "list number" in style_name:
                lines.append(f"1. {text}")
            else:
                lines.append(text)
        content = "\n".join(lines)
        return {"success": True, "content": content, "type": "docx", "path": str(p)}

    def _read_xlsx(self, p: Path) -> Dict[str, Any]:
        import openpyxl
        wb = openpyxl.load_workbook(str(p), data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## Feuille: {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    parts.append(" | ".join(cells))
        content = "\n".join(parts)
        return {"success": True, "content": content, "type": "xlsx", "path": str(p)}

    def _read_pptx(self, p: Path) -> Dict[str, Any]:
        from pptx import Presentation
        prs = Presentation(str(p))
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            parts.append(f"## Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
        content = "\n".join(parts)
        return {"success": True, "content": content, "type": "pptx", "path": str(p)}

    def _read_pdf(self, p: Path) -> Dict[str, Any]:
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(p))
            pages = []
            for i, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"--- Page {i} ---\n{text.strip()}")
            content = "\n\n".join(pages)
            return {"success": True, "content": content, "type": "pdf", "path": str(p)}
        except Exception as e:
            return {"success": False, "error": f"Impossible de lire le PDF: {e}"}

    # ──────────────────────────────────────────────────────
    # Supprimer un fichier
    # ──────────────────────────────────────────────────────

    def delete_file(self, path: str) -> Dict[str, Any]:
        """
        Supprime un fichier.
        Sécurité : seuls les fichiers dans le workspace peuvent être supprimés.
        """
        try:
            p = Path(path)
            if not p.is_absolute():
                found = self._find_in_workspace(path)
                if found is None:
                    return {"success": False, "error": f"Fichier non trouvé dans le workspace: {path}"}
                p = found

            if not p.exists():
                return {"success": False, "error": f"Fichier non trouvé: {p}"}

            # Vérification de sécurité : le fichier doit être dans le workspace
            try:
                p.resolve().relative_to(self.workspace_root.resolve())
            except ValueError:
                return {
                    "success": False,
                    "error": f"Sécurité: seuls les fichiers du workspace peuvent être supprimés. ({p})",
                }

            filename = p.name
            p.unlink()
            logger.info(f"🗑️ Fichier supprimé: {p}")
            return {"success": True, "path": str(p), "filename": filename}

        except Exception as e:
            logger.error(f"Erreur delete_file: {e}")
            return {"success": False, "error": str(e)}
# ──────────────────────────────────────────────────────────────────────────────
# © 2025-2026 LossKarr — Lumena Project
# Licensed under the GNU Affero General Public License v3.0 (AGPL-3.0)
# https://github.com/Losskarr/lumena
# ──────────────────────────────────────────────────────────────────────────────
