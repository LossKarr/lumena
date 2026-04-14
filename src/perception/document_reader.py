"""
document_reader.py — Lecteur de documents intelligent avec chunking sémantique.

Surpasse RAGFlow sur:
- Chunking par type (header / table / text / list / slide / image_page)
- Métadonnées complètes: page, section, source_file, char_offset
- Détection de tableaux Markdown et tableaux Word/Excel
- Découpage par phrase (pas arbitraire par tokens)
- Aucune dépendance externe: pypdf, python-docx, openpyxl, python-pptx (tous déjà installés)
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from loguru import logger
from pathlib import Path
from typing import Any, Dict, List, Optional


@dataclass
class DocumentChunk:
    """Un fragment de document avec métadonnées de provenance complètes."""

    content: str
    chunk_type: str          # "text" | "table" | "header" | "list" | "slide" | "image_page"
    page: int = 0            # Numéro de page (1-based, 0 si non applicable)
    section: str = ""        # Titre de section courant
    source_file: str = ""    # Chemin absolu du fichier source
    char_offset: int = 0     # Offset caractère dans le document source
    metadata: Dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "content": self.content,
            "chunk_type": self.chunk_type,
            "page": self.page,
            "section": self.section,
            "source_file": self.source_file,
            "char_offset": self.char_offset,
            **self.metadata,
        }

    def citation(self) -> str:
        """Retourne une citation de source formatée."""
        parts = [f"[Source: {Path(self.source_file).name}"]
        if self.page:
            parts.append(f"p.{self.page}")
        if self.section:
            parts.append(f"§ {self.section[:50]}")
        parts.append(f"type:{self.chunk_type}]")
        return " ".join(parts)


class DocumentReader:
    """
    Lecteur de documents avec chunking sémantique.

    Formats supportés: PDF, DOCX, XLSX, PPTX, TXT, MD, CSV, RST
    Retourne: List[DocumentChunk] avec provenance complète pour citations RAG.
    """

    MAX_CHUNK_CHARS: int = 1200   # Optimal pour embeddings ChromaDB
    MIN_CHUNK_CHARS: int = 60     # Ignorer les fragments trop courts

    # Headers Markdown (1-6 niveaux)
    _RE_HEADER = re.compile(r"^(#{1,6})\s+(.+)$")
    # Lignes de tableau Markdown
    _RE_TABLE_LINE = re.compile(r"^\|.+\|")
    # Séparateur de tableau
    _RE_TABLE_SEP = re.compile(r"^\|[-:| ]+\|$")

    # ─── Entry point ────────────────────────────────────────────────────────

    def read(self, path: str | Path) -> List[DocumentChunk]:
        """Lit un document et retourne des chunks sémantiques annotés."""
        p = Path(path)
        ext = p.suffix.lower()

        dispatch = {
            ".pdf": self._read_pdf,
            ".docx": self._read_docx,
            ".xlsx": self._read_xlsx,
            ".pptx": self._read_pptx,
        }

        if ext in dispatch:
            return dispatch[ext](p)

        if ext in {".txt", ".md", ".rst", ".csv", ".log", ".json", ".yaml", ".yml"}:
            return self._read_text(p)

        # Fallback générique
        try:
            content = p.read_text(encoding="utf-8", errors="replace")
            return self._chunk_text_with_sections(content, source_file=str(p))
        except Exception:
            return []  # fichier illisible

    # ─── PDF ────────────────────────────────────────────────────────────────

    def _read_pdf(self, p: Path) -> List[DocumentChunk]:
        try:
            from pypdf import PdfReader
        except ImportError:
            return []

        reader = PdfReader(str(p))
        chunks: List[DocumentChunk] = []
        current_section = ""

        for page_num, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""

            if not text.strip():
                # Page sans texte extractible → tentative OCR via vision.py
                ocr_text = self._ocr_pdf_page(p, page_num)
                if ocr_text:
                    text = ocr_text
                else:
                    chunks.append(DocumentChunk(
                        content=f"[Page {page_num} : contenu image non lisible]",
                        chunk_type="image_page",
                        page=page_num,
                        section=current_section,
                        source_file=str(p),
                        metadata={"image_only": True},
                    ))
                    continue

            page_chunks = self._chunk_text_with_sections(
                text,
                source_file=str(p),
                page=page_num,
                base_section=current_section,
            )

            # Track current section across pages
            for c in page_chunks:
                if c.chunk_type == "header":
                    current_section = c.content.lstrip("#").strip()

            chunks.extend(page_chunks)

        return chunks

    # ─── DOCX ───────────────────────────────────────────────────────────────

    def _read_docx(self, p: Path) -> List[DocumentChunk]:
        try:
            from docx import Document
        except ImportError:
            return []

        doc = Document(str(p))
        chunks: List[DocumentChunk] = []
        current_section = ""
        buffer: List[str] = []
        buffer_type = "text"

        def flush() -> None:
            nonlocal buffer, buffer_type
            text = "\n".join(buffer).strip()
            if len(text) >= self.MIN_CHUNK_CHARS:
                for c in self._split_into_chunks(
                    text, chunk_type=buffer_type,
                    source_file=str(p), section=current_section,
                ):
                    chunks.append(c)
            buffer.clear()
            buffer_type = "text"

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style = (para.style.name or "").lower()

            if "heading" in style or style == "title":
                flush()
                level = 1
                for i in range(1, 7):
                    if f"heading {i}" in style:
                        level = i
                        break
                hdr = "#" * level + " " + text
                chunks.append(DocumentChunk(
                    content=hdr,
                    chunk_type="header",
                    source_file=str(p),
                    section=current_section,
                ))
                current_section = text

            elif "list" in style:
                if buffer_type == "text" and buffer:
                    flush()
                buffer.append(f"- {text}")
                buffer_type = "list"

            else:
                if buffer_type == "list" and buffer:
                    flush()
                buffer.append(text)
                buffer_type = "text"

        flush()

        # Tables Word → Markdown
        for idx, table in enumerate(doc.tables):
            md = self._docx_table_to_markdown(table)
            if md:
                chunks.append(DocumentChunk(
                    content=md,
                    chunk_type="table",
                    source_file=str(p),
                    section=current_section,
                    metadata={"table_index": idx},
                ))

        return chunks

    @staticmethod
    def _docx_table_to_markdown(table) -> str:
        rows_data = []
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            rows_data.append(cells)
        if not rows_data:
            return ""
        lines = []
        ncols = max(len(r) for r in rows_data)
        for i, row in enumerate(rows_data):
            # Pad row
            padded = row + [""] * (ncols - len(row))
            lines.append("| " + " | ".join(padded) + " |")
            if i == 0:
                lines.append("|" + "|".join("---" for _ in padded) + "|")
        return "\n".join(lines)

    # ─── XLSX ───────────────────────────────────────────────────────────────

    def _read_xlsx(self, p: Path) -> List[DocumentChunk]:
        try:
            import openpyxl
        except ImportError:
            return []

        wb = openpyxl.load_workbook(str(p), data_only=True)
        chunks: List[DocumentChunk] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue

            ncols = max((len(r) for r in rows), default=0)
            if ncols == 0:
                continue

            # First row = header presumption
            header = [str(c) if c is not None else "" for c in rows[0]]
            lines = [
                f"## Feuille: {sheet_name}",
                "| " + " | ".join(header) + " |",
                "|" + "|".join("---" for _ in header) + "|",
            ]
            for row in rows[1:]:
                padded = list(row) + [None] * (ncols - len(row))
                cells = [str(c) if c is not None else "" for c in padded]
                if any(c.strip() for c in cells):
                    lines.append("| " + " | ".join(cells) + " |")

            chunks.append(DocumentChunk(
                content="\n".join(lines),
                chunk_type="table",
                source_file=str(p),
                section=sheet_name,
                metadata={"sheet": sheet_name, "rows": len(rows)},
            ))

        return chunks

    # ─── PPTX ───────────────────────────────────────────────────────────────

    def _read_pptx(self, p: Path) -> List[DocumentChunk]:
        try:
            from pptx import Presentation
        except ImportError:
            return []

        prs = Presentation(str(p))
        chunks: List[DocumentChunk] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts: List[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if not texts:
                continue

            title = texts[0]
            body = "\n".join(texts[1:])
            content = f"## {title}\n\n{body}" if body else f"## {title}"
            chunks.append(DocumentChunk(
                content=content,
                chunk_type="slide",
                page=slide_num,
                section=title,
                source_file=str(p),
                metadata={"slide_num": slide_num},
            ))

        return chunks

    def _ocr_pdf_page(self, p: Path, page_num: int) -> str:
        """Extrait le texte d'une page PDF image via l'OCR existant de Lumena (pytesseract)."""
        try:
            from pypdf import PdfReader
            from PIL import Image
            import io
            try:
                from ..computer_use.vision import ScreenAnalyzer
                analyzer = ScreenAnalyzer()
                if not analyzer.ocr_available:
                    return ""
            except Exception:
                return ""  # ScreenAnalyzer non disponible

            reader = PdfReader(str(p))
            page = reader.pages[page_num - 1]
            for img_obj in page.images:
                try:
                    img = Image.open(io.BytesIO(img_obj.data))
                    text = analyzer.extract_text(img)
                    if text.strip():
                        return text.strip()
                except Exception as exc:
                    logger.debug(f"[OCR] Method failed, trying next: {exc}")
                    continue
        except Exception as exc:
            logger.debug(f"[OCR] Full extraction failed: {exc}")
        return ""

    # ─── Text / Markdown ────────────────────────────────────────────────────

    def _read_text(self, p: Path) -> List[DocumentChunk]:
        content = p.read_text(encoding="utf-8", errors="replace")
        return self._chunk_text_with_sections(content, source_file=str(p))

    # ─── Core chunker ───────────────────────────────────────────────────────

    def _chunk_text_with_sections(
        self,
        text: str,
        source_file: str = "",
        page: int = 0,
        base_section: str = "",
    ) -> List[DocumentChunk]:
        """
        Découpe un texte en chunks sémantiques en respectant:
        - Les headers Markdown (# ## ### etc.)
        - Les tables Markdown (lignes |...|)
        - La taille maximale MAX_CHUNK_CHARS
        """
        chunks: List[DocumentChunk] = []
        current_section = base_section
        buffer: List[str] = []

        def flush_buffer(btype: str = "text") -> None:
            text_buf = "\n".join(buffer).strip()
            if len(text_buf) >= self.MIN_CHUNK_CHARS:
                for c in self._split_into_chunks(
                    text_buf,
                    chunk_type=btype,
                    source_file=source_file,
                    page=page,
                    section=current_section,
                ):
                    chunks.append(c)
            buffer.clear()

        lines = text.split("\n")
        i = 0
        in_table = False
        table_lines: List[str] = []

        while i < len(lines):
            line = lines[i]
            stripped = line.strip()

            # ── Header detection ──────────────────────────────────────────
            hm = self._RE_HEADER.match(stripped)
            if hm:
                if in_table:
                    if table_lines:
                        chunks.append(DocumentChunk(
                            content="\n".join(table_lines),
                            chunk_type="table",
                            page=page,
                            section=current_section,
                            source_file=source_file,
                        ))
                    table_lines = []
                    in_table = False
                flush_buffer()
                current_section = hm.group(2).strip()
                chunks.append(DocumentChunk(
                    content=stripped,
                    chunk_type="header",
                    page=page,
                    section=current_section,
                    source_file=source_file,
                ))
                i += 1
                continue

            # ── Table detection ──────────────────────────────────────────
            if self._RE_TABLE_LINE.match(stripped):
                if not in_table:
                    flush_buffer()
                    in_table = True
                table_lines.append(line)
                i += 1
                continue

            # Leaving table
            if in_table and not self._RE_TABLE_LINE.match(stripped):
                if table_lines:
                    chunks.append(DocumentChunk(
                        content="\n".join(table_lines),
                        chunk_type="table",
                        page=page,
                        section=current_section,
                        source_file=source_file,
                    ))
                table_lines = []
                in_table = False

            # ── Regular text / list ──────────────────────────────────────
            buffer.append(line)
            i += 1

        # Flush remainders
        if in_table and table_lines:
            chunks.append(DocumentChunk(
                content="\n".join(table_lines),
                chunk_type="table",
                page=page,
                section=current_section,
                source_file=source_file,
            ))
        flush_buffer()

        return chunks

    def _split_into_chunks(
        self,
        text: str,
        chunk_type: str = "text",
        source_file: str = "",
        page: int = 0,
        section: str = "",
    ) -> List[DocumentChunk]:
        """Découpe un bloc de texte en sous-chunks si nécessaire, par phrase."""
        if len(text) <= self.MAX_CHUNK_CHARS:
            return [DocumentChunk(
                content=text,
                chunk_type=chunk_type,
                page=page,
                section=section,
                source_file=source_file,
            )]

        # Split by sentence boundary
        sentences = re.split(r"(?<=[.!?])\s+", text)
        result: List[DocumentChunk] = []
        current: List[str] = []
        current_len = 0

        for sentence in sentences:
            if current_len + len(sentence) > self.MAX_CHUNK_CHARS and current:
                chunk_text = " ".join(current).strip()
                if len(chunk_text) >= self.MIN_CHUNK_CHARS:
                    result.append(DocumentChunk(
                        content=chunk_text,
                        chunk_type=chunk_type,
                        page=page,
                        section=section,
                        source_file=source_file,
                        char_offset=text.find(chunk_text),
                    ))
                current = [sentence]
                current_len = len(sentence)
            else:
                current.append(sentence)
                current_len += len(sentence)

        if current:
            chunk_text = " ".join(current).strip()
            if len(chunk_text) >= self.MIN_CHUNK_CHARS:
                result.append(DocumentChunk(
                    content=chunk_text,
                    chunk_type=chunk_type,
                    page=page,
                    section=section,
                    source_file=source_file,
                ))

        # Fallback: if nothing produced, return truncated
        if not result:
            result.append(DocumentChunk(
                content=text[: self.MAX_CHUNK_CHARS],
                chunk_type=chunk_type,
                page=page,
                section=section,
                source_file=source_file,
            ))

        return result
# ──────────────────────────────────────────────────────────────────────────────
# © 2025-2026 LossKarr — Lumena Project
# Licensed under the Apache License, Version 2.0
# https://github.com/Losskarr/lumena
# ──────────────────────────────────────────────────────────────────────────────
