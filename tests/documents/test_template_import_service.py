from __future__ import annotations

from pathlib import Path

from src.documents.studio import DocumentStudio


ROOT = Path(__file__).resolve().parents[2]


def _studio(tmp_path):
    return DocumentStudio(
        tmp_path / "studio", builtin_root=ROOT / "assets" / "templates",
        output_root=tmp_path / "output",
    )


def test_html_import_is_a_reviewable_draft_then_a_routable_custom_template(tmp_path):
    studio = _studio(tmp_path)
    source = tmp_path / "modele.html"
    source.write_text(
        "<html><body onload='evil()'><script>alert(1)</script>"
        "<h1>Rapport [[client_name]]</h1><p>Période ${period}</p></body></html>",
        encoding="utf-8",
    )

    draft = studio.template_imports.create(
        source, filename=source.name, name="Rapport client", kind="rapport_client",
    )

    assert draft.status == "draft"
    assert draft.fidelity == "high"
    assert {field["id"] for field in draft.detected_fields} == {"client_name", "period"}
    assert "<script" not in draft.template_source
    assert "onload" not in draft.template_source
    assert draft.sample_data == {"client_name": "client_name", "period": "period"}
    assert studio.catalog.get_default("rapport_client", "pdf") is None

    record = studio.template_imports.publish(draft.id, template_id="rapport-client-atlas")
    assert record.manifest.id == "rapport-client-atlas"
    assert record.manifest.kind == "rapport_client"
    assert studio.catalog.get("rapport-client-atlas").read_only is False
    assert studio.template_imports.get(draft.id).status == "published"


def test_docx_import_preserves_paragraphs_tables_and_placeholders(tmp_path):
    from docx import Document

    studio = _studio(tmp_path)
    source = tmp_path / "contrat.docx"
    document = Document()
    document.add_heading("Contrat {{client}}", 0)
    document.add_paragraph("Référence [[reference]]")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Montant"
    table.cell(0, 1).text = "${amount}"
    document.save(source)

    draft = studio.template_imports.create(source, kind="contrat_importe")

    assert draft.source_format == "docx"
    assert draft.fidelity == "structural"
    assert "<table>" in draft.template_source
    assert {field["id"] for field in draft.detected_fields} == {"client", "reference", "amount"}
    assert "{{ client }}" in studio.template_imports.preview_html(draft.id) or "client" in draft.sample_data


def test_xlsx_and_pptx_sources_create_structural_drafts(tmp_path):
    from openpyxl import Workbook
    from pptx import Presentation

    studio = _studio(tmp_path)
    workbook = Workbook()
    workbook.active.title = "Budget"
    workbook.active.append(["Client", "[[client]]"])
    xlsx = tmp_path / "budget.xlsx"
    workbook.save(xlsx)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Présentation {{company}}"
    pptx = tmp_path / "pitch.pptx"
    presentation.save(pptx)

    sheet_draft = studio.template_imports.create(xlsx, kind="budget_importe")
    slide_draft = studio.template_imports.create(pptx, kind="pitch_importe")

    assert sheet_draft.source_format == "xlsx"
    assert sheet_draft.detected_fields[0]["id"] == "client"
    assert "<table>" in sheet_draft.template_source
    assert slide_draft.source_format == "pptx"
    assert slide_draft.detected_fields[0]["id"] == "company"
    assert "imported-slide" in slide_draft.template_source


def test_source_without_placeholders_stays_draft_and_warns(tmp_path):
    studio = _studio(tmp_path)
    source = tmp_path / "reference.txt"
    source.write_text("Document fixe sans champ variable", encoding="utf-8")

    draft = studio.template_imports.create(source, kind="reference_fixe")

    assert draft.status == "draft"
    assert draft.detected_fields == ()
    assert "no_placeholders_detected" in draft.warnings


def test_published_draft_cannot_be_deleted_or_republished(tmp_path):
    import pytest

    studio = _studio(tmp_path)
    source = tmp_path / "simple.html"
    source.write_text("<html><body>{{ value }}</body></html>", encoding="utf-8")
    draft = studio.template_imports.create(source, kind="simple_modele")
    studio.template_imports.publish(draft.id)

    with pytest.raises(ValueError, match="déjà publié"):
        studio.template_imports.publish(draft.id)
    with pytest.raises(ValueError, match="preuve de provenance"):
        studio.template_imports.delete(draft.id)
