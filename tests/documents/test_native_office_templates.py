from __future__ import annotations

from pathlib import Path

import pytest

from src.documents.document_validation import validate_document_render
from src.documents.studio import DocumentStudio
from src.documents.template_models import TemplateValidationError


ROOT = Path(__file__).resolve().parents[2]


def _studio(tmp_path: Path) -> DocumentStudio:
    return DocumentStudio(
        tmp_path / "studio",
        builtin_root=ROOT / "assets" / "templates",
        output_root=tmp_path / "output",
    )


@pytest.mark.asyncio
async def test_docx_import_publishes_and_generates_a_native_document(tmp_path):
    from docx import Document

    studio = _studio(tmp_path)
    source = tmp_path / "offer.docx"
    document = Document()
    paragraph = document.add_paragraph()
    paragraph.add_run("Offer {{ cli")
    paragraph.add_run("ent }}")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Reference"
    table.cell(0, 1).text = "[[reference]]"
    document.save(source)

    draft = studio.template_imports.create(source, kind="native_offer")
    assert (draft.renderer, draft.output_format) == ("docx-native", "docx")
    record = studio.template_imports.publish(draft.id)
    assert record.manifest.template_file == "template.docx"
    with pytest.raises(TemplateValidationError, match="do not expose"):
        studio.catalog.read_source(record)

    result = await studio.generate(
        template_id=record.manifest.id,
        data={"client": "Atlas", "reference": "OFF-42"},
        filename="offer-atlas",
    )
    output = Path(result["path"])
    rendered = Document(output)
    assert output.suffix == ".docx"
    assert rendered.paragraphs[0].text == "Offer Atlas"
    assert rendered.tables[0].cell(0, 1).text == "OFF-42"
    assert result["render_proof"]["status"] == "structure_verified"
    assert result["render_proof"]["verification_level"] == "structural"


@pytest.mark.asyncio
async def test_xlsx_native_render_preserves_formulas_and_replaces_cells(tmp_path):
    from openpyxl import Workbook, load_workbook

    studio = _studio(tmp_path)
    source = tmp_path / "budget.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.append(["Client", "[[client]]"])
    sheet.append([2, 3])
    sheet["C2"] = "=A2+B2"
    workbook.save(source)
    workbook.close()

    draft = studio.template_imports.create(source, kind="native_budget")
    record = studio.template_imports.publish(draft.id)
    result = await studio.generate(
        template_id=record.manifest.id,
        data={"client": "Nova"},
        filename="budget-nova",
    )
    rendered = load_workbook(result["path"], data_only=False)
    try:
        assert rendered.active["B1"].value == "Nova"
        assert rendered.active["C2"].value == "=A2+B2"
    finally:
        rendered.close()
    assert result["render_proof"]["status"] == "structure_verified"


@pytest.mark.asyncio
async def test_pptx_native_render_replaces_text_and_keeps_slide(tmp_path):
    from pptx import Presentation

    studio = _studio(tmp_path)
    source = tmp_path / "pitch.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Pitch <<company>>"
    presentation.save(source)

    draft = studio.template_imports.create(source, kind="native_pitch")
    record = studio.template_imports.publish(draft.id)
    result = await studio.generate(
        template_id=record.manifest.id,
        data={"company": "Lumena"},
        filename="pitch-lumena",
    )
    rendered = Presentation(result["path"])
    assert len(rendered.slides) == 1
    assert rendered.slides[0].shapes.title.text == "Pitch Lumena"
    assert result["render_proof"]["page_count"] == 1


def test_native_validation_is_structural_without_claiming_visual_proof(tmp_path):
    from docx import Document

    source = tmp_path / "proof.docx"
    document = Document()
    document.add_paragraph("Non-empty native document")
    document.save(source)

    proof = validate_document_render(
        source,
        None,
        output_format="docx",
        visual_fidelity="structural",
    )
    assert proof["verified"] is True
    assert proof["status"] == "structure_verified"
    assert proof["verification_level"] == "structural"
    assert proof["thumbnail_path"] == ""


def test_native_template_rejects_missing_data_instead_of_leaving_placeholder(tmp_path):
    from docx import Document

    studio = _studio(tmp_path)
    source = tmp_path / "strict.docx"
    document = Document()
    document.add_paragraph("Client {{ client }}")
    document.save(source)
    draft = studio.template_imports.create(source, kind="native_strict")
    record = studio.template_imports.publish(draft.id)

    with pytest.raises(TemplateValidationError, match="missing render field"):
        studio.renderer.render_native(
            record,
            studio.catalog.source_path(record),
            {},
            tmp_path / "result.docx",
        )
